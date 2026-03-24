"""Tests for doc_export.py — PPTX, DOCX, and PDF generation from lesson plans."""

from __future__ import annotations

from pathlib import Path
from unittest.mock import patch

import pytest

from eduagent.models import (
    DailyLesson,
    DifferentiationNotes,
    ExitTicketQuestion,
    TeacherPersona,
)

# ── Fixtures ──────────────────────────────────────────────────────────


@pytest.fixture
def sample_lesson() -> DailyLesson:
    """A realistic DailyLesson with all required fields populated."""
    return DailyLesson(
        title="The Causes of World War I",
        lesson_number=3,
        objective="Students will be able to identify the four main causes of WWI (MAIN).",
        standards=["D2.His.1.9-12", "D2.His.3.9-12"],
        do_now="List one thing you already know about World War I.",
        direct_instruction=(
            "World War I was triggered by a complex web of factors often summarized "
            "by the acronym MAIN: Militarism, Alliances, Imperialism, and Nationalism. "
            "Each of these forces contributed to rising tensions across Europe in the "
            "early 20th century."
        ),
        guided_practice=(
            "In pairs, analyze a political cartoon from 1914. Identify which MAIN "
            "factor the cartoon illustrates and explain your reasoning."
        ),
        independent_work=(
            "Complete the MAIN graphic organizer by writing a one-sentence summary "
            "and a real-world example for each factor."
        ),
        exit_ticket=[
            ExitTicketQuestion(
                question="Which MAIN factor do you think was most responsible for WWI? Why?",
                expected_response="Any of the four with supporting reasoning.",
            ),
        ],
        homework="Read pages 112-118 and take notes on the assassination of Archduke Franz Ferdinand.",
        differentiation=DifferentiationNotes(
            struggling=["Provide a pre-filled graphic organizer with word bank."],
            advanced=["Compare MAIN factors to causes of a modern conflict."],
            ell=["Provide translated key terms in Spanish."],
        ),
        materials_needed=["MAIN graphic organizer", "Political cartoon handout", "Textbook"],
        time_estimates={
            "do_now": 5,
            "direct_instruction": 15,
            "guided_practice": 15,
            "independent_work": 10,
        },
    )


@pytest.fixture
def sample_persona() -> TeacherPersona:
    """A simple teacher persona for export tests."""
    return TeacherPersona(
        name="Ms. Rivera",
        subject_area="Social Studies",
        grade_levels=["8"],
    )


@pytest.fixture
def minimal_lesson() -> DailyLesson:
    """A DailyLesson with only required fields and minimal optional data."""
    return DailyLesson(
        title="Quick Lesson",
        lesson_number=1,
        objective="SWBAT understand the basics.",
    )


@pytest.fixture
def long_title_lesson() -> DailyLesson:
    """A DailyLesson with a very long title for edge-case testing."""
    return DailyLesson(
        title="A" * 200,
        lesson_number=99,
        objective="Test extremely long title handling.",
        do_now="Warm-up activity.",
        direct_instruction="Instruction content " * 50,
        guided_practice="Practice content.",
        independent_work="Work content.",
    )


# ── PPTX export ──────────────────────────────────────────────────────


class TestPPTXExport:
    def test_generates_pptx_file(self, sample_lesson, sample_persona, tmp_path):
        from eduagent.doc_export import export_lesson_pptx

        path = export_lesson_pptx(sample_lesson, sample_persona, output_dir=tmp_path)
        assert path.exists()
        assert path.suffix == ".pptx"

    def test_pptx_is_nonempty(self, sample_lesson, sample_persona, tmp_path):
        from eduagent.doc_export import export_lesson_pptx

        path = export_lesson_pptx(sample_lesson, sample_persona, output_dir=tmp_path)
        assert path.stat().st_size > 0

    def test_pptx_has_correct_slide_count(self, sample_lesson, sample_persona, tmp_path):
        from pptx import Presentation

        from eduagent.doc_export import export_lesson_pptx

        path = export_lesson_pptx(sample_lesson, sample_persona, output_dir=tmp_path)
        prs = Presentation(str(path))
        # Title + Objectives + DoNow + DirectInstruction + GuidedPractice
        # + IndependentWork + ExitTicket + Homework = 8 slides
        assert len(prs.slides) >= 5  # At minimum title + objectives + 3 sections

    def test_pptx_minimal_lesson(self, minimal_lesson, sample_persona, tmp_path):
        from eduagent.doc_export import export_lesson_pptx

        path = export_lesson_pptx(minimal_lesson, sample_persona, output_dir=tmp_path)
        assert path.exists()
        assert path.stat().st_size > 0

    def test_pptx_long_title(self, long_title_lesson, sample_persona, tmp_path):
        from eduagent.doc_export import export_lesson_pptx

        path = export_lesson_pptx(long_title_lesson, sample_persona, output_dir=tmp_path)
        assert path.exists()
        assert path.stat().st_size > 0


# ── DOCX export ──────────────────────────────────────────────────────


class TestDOCXExport:
    def test_generates_docx_file(self, sample_lesson, sample_persona, tmp_path):
        from eduagent.doc_export import export_lesson_docx

        path = export_lesson_docx(sample_lesson, sample_persona, output_dir=tmp_path)
        assert path.exists()
        assert path.suffix == ".docx"

    def test_docx_is_nonempty(self, sample_lesson, sample_persona, tmp_path):
        from eduagent.doc_export import export_lesson_docx

        path = export_lesson_docx(sample_lesson, sample_persona, output_dir=tmp_path)
        assert path.stat().st_size > 0

    def test_docx_contains_lesson_title(self, sample_lesson, sample_persona, tmp_path):
        from docx import Document

        from eduagent.doc_export import export_lesson_docx

        path = export_lesson_docx(sample_lesson, sample_persona, output_dir=tmp_path)
        doc = Document(str(path))
        full_text = "\n".join(p.text for p in doc.paragraphs)
        assert "Causes of World War I" in full_text

    def test_docx_minimal_lesson(self, minimal_lesson, sample_persona, tmp_path):
        from eduagent.doc_export import export_lesson_docx

        path = export_lesson_docx(minimal_lesson, sample_persona, output_dir=tmp_path)
        assert path.exists()
        assert path.stat().st_size > 0

    def test_docx_long_title(self, long_title_lesson, sample_persona, tmp_path):
        from eduagent.doc_export import export_lesson_docx

        path = export_lesson_docx(long_title_lesson, sample_persona, output_dir=tmp_path)
        assert path.exists()
        assert path.stat().st_size > 0


# ── PDF export ────────────────────────────────────────────────────────


class TestPDFExport:
    def test_generates_pdf_file(self, sample_lesson, sample_persona, tmp_path):
        from eduagent.doc_export import export_lesson_pdf

        path = export_lesson_pdf(sample_lesson, sample_persona, output_dir=tmp_path)
        assert path.exists()
        assert path.suffix == ".pdf"

    def test_pdf_is_nonempty(self, sample_lesson, sample_persona, tmp_path):
        from eduagent.doc_export import export_lesson_pdf

        path = export_lesson_pdf(sample_lesson, sample_persona, output_dir=tmp_path)
        assert path.stat().st_size > 0

    def test_pdf_minimal_lesson(self, minimal_lesson, sample_persona, tmp_path):
        from eduagent.doc_export import export_lesson_pdf

        path = export_lesson_pdf(minimal_lesson, sample_persona, output_dir=tmp_path)
        assert path.exists()
        assert path.stat().st_size > 0

    def test_pdf_long_title(self, long_title_lesson, sample_persona, tmp_path):
        from eduagent.doc_export import export_lesson_pdf

        path = export_lesson_pdf(long_title_lesson, sample_persona, output_dir=tmp_path)
        assert path.exists()
        assert path.stat().st_size > 0


# ── Color theme tests ────────────────────────────────────────────────


class TestColorTheme:
    def test_history_theme(self):
        from eduagent.doc_export import get_color_theme

        theme = get_color_theme("history")
        assert theme["primary"] == "8B4513"  # Saddle brown
        assert theme["secondary"] == "DAA520"  # Goldenrod

    def test_social_studies_matches_history(self):
        from eduagent.doc_export import get_color_theme

        assert get_color_theme("social studies") == get_color_theme("history")

    def test_science_theme(self):
        from eduagent.doc_export import get_color_theme

        theme = get_color_theme("science")
        assert theme["primary"] == "1B5E20"  # Dark green

    def test_math_theme(self):
        from eduagent.doc_export import get_color_theme

        theme = get_color_theme("math")
        assert theme["primary"] == "1565C0"  # Blue

    def test_ela_theme(self):
        from eduagent.doc_export import get_color_theme

        theme = get_color_theme("ela")
        assert theme["primary"] == "6A1B9A"  # Purple

    def test_default_theme_for_unknown_subject(self):
        from eduagent.doc_export import get_color_theme

        theme = get_color_theme("underwater basket weaving")
        assert theme["primary"] == "1A365D"  # professional navy

    def test_case_insensitive(self):
        from eduagent.doc_export import get_color_theme

        assert get_color_theme("History") == get_color_theme("history")
        assert get_color_theme("MATH") == get_color_theme("math")

    def test_theme_has_all_required_keys(self):
        from eduagent.doc_export import get_color_theme

        for subject in ["history", "science", "math", "ela", "unknown"]:
            theme = get_color_theme(subject)
            for key in ("primary", "secondary", "accent", "bg_dark", "bg_light", "text_dark", "text_light"):
                assert key in theme, f"Missing key '{key}' in theme for '{subject}'"


# ── Search query building tests ──────────────────────────────────────


class TestSearchQueryBuilding:
    def test_basic_topic(self):
        from eduagent.slide_images import _build_search_query

        query = _build_search_query("Photosynthesis", "science")
        assert "photosynthesis" in query
        # Science subject adds "diagram" or "illustration"
        assert "diagram" in query or "illustration" in query

    def test_strips_articles(self):
        from eduagent.slide_images import _build_search_query

        query = _build_search_query("The American Revolution", "history")
        assert "the" not in query.split()
        assert "american" in query
        assert "revolution" in query

    def test_history_adds_primary_source(self):
        from eduagent.slide_images import _build_search_query

        query = _build_search_query("The Civil War", "history")
        assert "historical" in query or "primary source" in query

    def test_adds_education_keywords(self):
        from eduagent.slide_images import _build_search_query

        query = _build_search_query("Solving Linear Equations", "math")
        assert "mathematics" in query or "education" in query

    def test_no_subject_uses_education_default(self):
        from eduagent.slide_images import _build_search_query

        query = _build_search_query("Random Topic")
        assert "education" in query

    def test_no_duplicate_words(self):
        from eduagent.slide_images import _build_search_query

        query = _build_search_query("mathematics quiz", "math")
        words = query.split()
        assert words.count("mathematics") <= 1


# ── PPTX with mock images ───────────────────────────────────────────


class TestPPTXWithImages:
    @pytest.fixture
    def fake_image(self, tmp_path):
        """Create a minimal valid JPEG file for testing."""
        # Minimal JPEG: SOI marker + bare minimum to be recognized
        # We'll use a tiny valid PNG instead (simpler to construct)
        import struct
        import zlib

        def _create_png(path: Path):
            """Create a 1x1 white PNG file."""
            sig = b"\x89PNG\r\n\x1a\n"
            # IHDR
            ihdr_data = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
            ihdr_crc = zlib.crc32(b"IHDR" + ihdr_data) & 0xFFFFFFFF
            ihdr = struct.pack(">I", 13) + b"IHDR" + ihdr_data + struct.pack(">I", ihdr_crc)
            # IDAT
            raw = b"\x00\xff\xff\xff"  # filter=none, R=255, G=255, B=255
            compressed = zlib.compress(raw)
            idat_crc = zlib.crc32(b"IDAT" + compressed) & 0xFFFFFFFF
            idat = struct.pack(">I", len(compressed)) + b"IDAT" + compressed + struct.pack(">I", idat_crc)
            # IEND
            iend_crc = zlib.crc32(b"IEND") & 0xFFFFFFFF
            iend = struct.pack(">I", 0) + b"IEND" + struct.pack(">I", iend_crc)
            path.write_bytes(sig + ihdr + idat + iend)

        img_path = tmp_path / "test_image.png"
        _create_png(img_path)
        return img_path

    def test_pptx_with_mock_image(self, sample_lesson, sample_persona, tmp_path, fake_image):
        """PPTX export works when fetch_slide_image returns a real image."""
        from eduagent.doc_export import export_lesson_pptx

        async def mock_fetch(topic, subject="", cache_dir=None):
            return fake_image

        with patch("eduagent.slide_images.fetch_slide_image", side_effect=mock_fetch):
            path = export_lesson_pptx(sample_lesson, sample_persona, output_dir=tmp_path)

        assert path.exists()
        assert path.stat().st_size > 0

        from pptx import Presentation
        prs = Presentation(str(path))
        assert len(prs.slides) >= 5

    def test_pptx_without_images_api_unavailable(self, sample_lesson, sample_persona, tmp_path):
        """PPTX export works fine when no Unsplash key is configured."""
        from eduagent.doc_export import export_lesson_pptx

        # Ensure no UNSPLASH_ACCESS_KEY is set
        with patch.dict("os.environ", {}, clear=False):
            with patch("eduagent.slide_images._get_unsplash_key", return_value=None):
                path = export_lesson_pptx(sample_lesson, sample_persona, output_dir=tmp_path)

        assert path.exists()
        assert path.stat().st_size > 0

        from pptx import Presentation
        prs = Presentation(str(path))
        # Should still have all slides even without images
        assert len(prs.slides) >= 5

    def test_pptx_graceful_on_fetch_failure(self, sample_lesson, sample_persona, tmp_path):
        """PPTX export completes even when image fetching raises exceptions."""
        from eduagent.doc_export import export_lesson_pptx

        async def failing_fetch(topic, subject="", cache_dir=None):
            raise ConnectionError("Network unavailable")

        with patch("eduagent.slide_images.fetch_slide_image", side_effect=failing_fetch):
            with patch("eduagent.slide_images._get_unsplash_key", return_value="fake-key"):
                path = export_lesson_pptx(sample_lesson, sample_persona, output_dir=tmp_path)

        assert path.exists()
        assert path.stat().st_size > 0


# ── PPTX theme integration tests ────────────────────────────────────


class TestPPTXThemeIntegration:
    def test_history_persona_uses_warm_theme(self, sample_lesson, tmp_path):
        """A Social Studies persona produces slides with warm-toned theme."""
        from pptx import Presentation

        from eduagent.doc_export import export_lesson_pptx

        persona = TeacherPersona(name="Mr. Smith", subject_area="History")
        path = export_lesson_pptx(sample_lesson, persona, output_dir=tmp_path)
        assert path.exists()

        prs = Presentation(str(path))
        assert len(prs.slides) >= 5

    def test_science_persona_uses_cool_theme(self, sample_lesson, tmp_path):
        from eduagent.doc_export import export_lesson_pptx

        persona = TeacherPersona(name="Dr. Lee", subject_area="Science")
        path = export_lesson_pptx(sample_lesson, persona, output_dir=tmp_path)
        assert path.exists()
        assert path.stat().st_size > 0

    def test_math_persona_uses_blue_theme(self, sample_lesson, tmp_path):
        from eduagent.doc_export import export_lesson_pptx

        persona = TeacherPersona(name="Ms. Chen", subject_area="Math")
        path = export_lesson_pptx(sample_lesson, persona, output_dir=tmp_path)
        assert path.exists()

    def test_ela_persona_uses_purple_theme(self, sample_lesson, tmp_path):
        from eduagent.doc_export import export_lesson_pptx

        persona = TeacherPersona(name="Mrs. Jones", subject_area="ELA")
        path = export_lesson_pptx(sample_lesson, persona, output_dir=tmp_path)
        assert path.exists()

    def test_unknown_subject_uses_default(self, sample_lesson, tmp_path):
        from eduagent.doc_export import export_lesson_pptx

        persona = TeacherPersona(name="Coach", subject_area="Physical Education")
        path = export_lesson_pptx(sample_lesson, persona, output_dir=tmp_path)
        assert path.exists()

    def test_closing_slide_has_homework(self, sample_lesson, sample_persona, tmp_path):
        """When lesson has homework, closing slide should include it."""
        from pptx import Presentation

        from eduagent.doc_export import export_lesson_pptx

        path = export_lesson_pptx(sample_lesson, sample_persona, output_dir=tmp_path)
        prs = Presentation(str(path))
        last_slide = prs.slides[-1]
        texts = [shape.text for shape in last_slide.shapes if shape.has_text_frame]
        combined = " ".join(texts)
        assert "Homework" in combined or "homework" in combined.lower()

    def test_closing_slide_questions_when_no_homework(self, sample_persona, tmp_path):
        """When no homework, closing slide should say 'Questions?'."""
        from pptx import Presentation

        from eduagent.doc_export import export_lesson_pptx

        lesson = DailyLesson(
            title="No Homework Lesson",
            lesson_number=1,
            objective="Test closing slide.",
            homework=None,
        )
        path = export_lesson_pptx(lesson, sample_persona, output_dir=tmp_path)
        prs = Presentation(str(path))
        last_slide = prs.slides[-1]
        texts = [shape.text for shape in last_slide.shapes if shape.has_text_frame]
        combined = " ".join(texts)
        assert "Questions?" in combined


# ── Academic image source selection tests ───────────────────────────


class TestAcademicImageSources:
    """Tests for the multi-source image routing in slide_images.py."""

    def test_history_prefers_loc(self):
        from eduagent.slide_images import _select_sources

        sources = _select_sources("History")
        assert sources[0] == "loc"
        assert "wikimedia" in sources
        assert "unsplash" in sources

    def test_social_studies_prefers_loc(self):
        from eduagent.slide_images import _select_sources

        sources = _select_sources("Social Studies")
        assert sources[0] == "loc"

    def test_civics_prefers_loc(self):
        from eduagent.slide_images import _select_sources

        sources = _select_sources("Civics")
        assert sources[0] == "loc"

    def test_science_prefers_wikimedia(self):
        from eduagent.slide_images import _select_sources

        sources = _select_sources("Science")
        assert sources[0] == "wikimedia"

    def test_biology_prefers_wikimedia(self):
        from eduagent.slide_images import _select_sources

        sources = _select_sources("Biology")
        assert sources[0] == "wikimedia"

    def test_art_prefers_wikimedia(self):
        from eduagent.slide_images import _select_sources

        sources = _select_sources("Art")
        assert sources[0] == "wikimedia"

    def test_math_prefers_unsplash(self):
        from eduagent.slide_images import _select_sources

        sources = _select_sources("Math")
        assert sources[0] == "unsplash"
        assert "wikimedia" in sources

    def test_unknown_subject_defaults_to_unsplash_first(self):
        from eduagent.slide_images import _select_sources

        sources = _select_sources("Underwater Basket Weaving")
        assert sources[0] == "unsplash"


class TestLOCResponseParsing:
    """Test parsing of Library of Congress API responses."""

    def test_parses_image_url_list(self):
        """LOC results with image_url list should return the first URL."""
        import asyncio
        from unittest.mock import AsyncMock, MagicMock, patch

        from eduagent.slide_images import _fetch_loc

        mock_response = MagicMock()
        mock_response.status_code = 200
        mock_response.raise_for_status = MagicMock()
        mock_response.json.return_value = {
            "results": [
                {
                    "image_url": [
                        "https://tile.loc.gov/image-services/iiif/test.jpg",
                    ],
                },
            ],
        }

        mock_img_resp = MagicMock()
        mock_img_resp.raise_for_status = MagicMock()
        mock_img_resp.content = b"\xff\xd8\xff\xe0fake-jpeg"

        async def mock_get(url, **kwargs):
            if "loc.gov/search" in str(url):
                return mock_response
            return mock_img_resp

        with patch("httpx.AsyncClient") as mock_client_cls:
            mock_client = AsyncMock()
            mock_client.get = mock_get
            mock_client.__aenter__ = AsyncMock(return_value=mock_client)
            mock_client.__aexit__ = AsyncMock(return_value=False)
            mock_client_cls.return_value = mock_client

            result = asyncio.run(_fetch_loc("civil war", cache_dir=None))
            # Should attempt to fetch -- result depends on whether cache wrote
            # but should not raise
            assert result is None or isinstance(result, Path)

    def test_handles_empty_results(self):
        import asyncio
        from unittest.mock import AsyncMock, MagicMock, patch

        from eduagent.slide_images import _fetch_loc

        mock_response = MagicMock()
        mock_response.status_code = 200
        mock_response.raise_for_status = MagicMock()
        mock_response.json.return_value = {"results": []}

        async def mock_get(url, **kwargs):
            return mock_response

        with patch("httpx.AsyncClient") as mock_client_cls:
            mock_client = AsyncMock()
            mock_client.get = mock_get
            mock_client.__aenter__ = AsyncMock(return_value=mock_client)
            mock_client.__aexit__ = AsyncMock(return_value=False)
            mock_client_cls.return_value = mock_client

            result = asyncio.run(_fetch_loc("nonexistent query xyz"))
            assert result is None


class TestWikimediaResponseParsing:
    """Test parsing of Wikimedia Commons API responses."""

    def test_handles_empty_response(self):
        import asyncio
        from unittest.mock import AsyncMock, MagicMock, patch

        from eduagent.slide_images import _fetch_wikimedia

        mock_response = MagicMock()
        mock_response.status_code = 200
        mock_response.raise_for_status = MagicMock()
        mock_response.json.return_value = {"query": {"pages": {}}}

        async def mock_get(url, **kwargs):
            return mock_response

        with patch("httpx.AsyncClient") as mock_client_cls:
            mock_client = AsyncMock()
            mock_client.get = mock_get
            mock_client.__aenter__ = AsyncMock(return_value=mock_client)
            mock_client.__aexit__ = AsyncMock(return_value=False)
            mock_client_cls.return_value = mock_client

            result = asyncio.run(_fetch_wikimedia("nonexistent query xyz"))
            assert result is None

    def test_handles_network_error(self):
        import asyncio
        from unittest.mock import AsyncMock, patch

        from eduagent.slide_images import _fetch_wikimedia

        async def mock_get(url, **kwargs):
            raise ConnectionError("Network down")

        with patch("httpx.AsyncClient") as mock_client_cls:
            mock_client = AsyncMock()
            mock_client.get = mock_get
            mock_client.__aenter__ = AsyncMock(return_value=mock_client)
            mock_client.__aexit__ = AsyncMock(return_value=False)
            mock_client_cls.return_value = mock_client

            result = asyncio.run(_fetch_wikimedia("photosynthesis"))
            assert result is None


class TestSlideCountMatchesContent:
    """Verify the slide count matches what the lesson content should produce."""

    def test_full_lesson_slide_count(self, sample_lesson, sample_persona, tmp_path):
        """A lesson with all sections should produce 8 slides."""
        from pptx import Presentation

        from eduagent.doc_export import export_lesson_pptx

        path = export_lesson_pptx(sample_lesson, sample_persona, output_dir=tmp_path)
        prs = Presentation(str(path))
        # Title + Objectives + DoNow + DirectInstruction + GuidedPractice
        # + IndependentWork + ExitTicket + Closing = 8
        assert len(prs.slides) == 8

    def test_minimal_lesson_slide_count(self, minimal_lesson, sample_persona, tmp_path):
        """A minimal lesson (title + objective only) should produce 3 slides."""
        from pptx import Presentation

        from eduagent.doc_export import export_lesson_pptx

        path = export_lesson_pptx(minimal_lesson, sample_persona, output_dir=tmp_path)
        prs = Presentation(str(path))
        # Title + Objectives + Closing = 3
        assert len(prs.slides) == 3

    def test_no_exit_ticket_reduces_count(self, sample_persona, tmp_path):
        """Removing exit ticket should reduce slide count by 1."""
        from pptx import Presentation

        from eduagent.doc_export import export_lesson_pptx

        lesson = DailyLesson(
            title="No Exit Ticket Lesson",
            lesson_number=1,
            objective="Test slide count.",
            do_now="Quick warm-up.",
            direct_instruction="Main content.",
            guided_practice="Practice activity.",
            independent_work="Solo work.",
            exit_ticket=[],
            homework="Read chapter 5.",
        )
        path = export_lesson_pptx(lesson, sample_persona, output_dir=tmp_path)
        prs = Presentation(str(path))
        # Title + Objectives + DoNow + DI + GP + IW + Closing = 7
        assert len(prs.slides) == 7


# ── Student handout export ─────────────────────────────────────────


class TestStudentHandoutExport:
    def test_generates_handout_file(self, sample_lesson, sample_persona, tmp_path):
        from eduagent.doc_export import export_student_handout

        path = export_student_handout(sample_lesson, sample_persona, output_dir=tmp_path)
        assert path.exists()
        assert path.suffix == ".docx"
        assert "_handout" in path.name

    def test_handout_is_nonempty(self, sample_lesson, sample_persona, tmp_path):
        from eduagent.doc_export import export_student_handout

        path = export_student_handout(sample_lesson, sample_persona, output_dir=tmp_path)
        assert path.stat().st_size > 0

    def test_handout_contains_lesson_title(self, sample_lesson, sample_persona, tmp_path):
        from docx import Document

        from eduagent.doc_export import export_student_handout

        path = export_student_handout(sample_lesson, sample_persona, output_dir=tmp_path)
        doc = Document(str(path))
        full_text = "\n".join(p.text for p in doc.paragraphs)
        assert "Causes of World War I" in full_text

    def test_handout_contains_objective(self, sample_lesson, sample_persona, tmp_path):
        from docx import Document

        from eduagent.doc_export import export_student_handout

        path = export_student_handout(sample_lesson, sample_persona, output_dir=tmp_path)
        doc = Document(str(path))
        full_text = "\n".join(p.text for p in doc.paragraphs)
        assert "MAIN" in full_text

    def test_handout_contains_do_now(self, sample_lesson, sample_persona, tmp_path):
        from docx import Document

        from eduagent.doc_export import export_student_handout

        path = export_student_handout(sample_lesson, sample_persona, output_dir=tmp_path)
        doc = Document(str(path))
        # Do now should be in a table cell
        all_text = ""
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    all_text += cell.text + "\n"
        assert "World War I" in all_text or "already know" in all_text

    def test_handout_contains_exit_ticket(self, sample_lesson, sample_persona, tmp_path):
        from docx import Document

        from eduagent.doc_export import export_student_handout

        path = export_student_handout(sample_lesson, sample_persona, output_dir=tmp_path)
        doc = Document(str(path))
        all_text = ""
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    all_text += cell.text + "\n"
        assert "MAIN factor" in all_text

    def test_handout_contains_footer_fields(self, sample_lesson, sample_persona, tmp_path):
        from docx import Document

        from eduagent.doc_export import export_student_handout

        path = export_student_handout(sample_lesson, sample_persona, output_dir=tmp_path)
        doc = Document(str(path))
        all_text = ""
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    all_text += cell.text + "\n"
        assert "Name:" in all_text
        assert "Date:" in all_text
        assert "Period:" in all_text

    def test_handout_minimal_lesson(self, minimal_lesson, sample_persona, tmp_path):
        from eduagent.doc_export import export_student_handout

        path = export_student_handout(minimal_lesson, sample_persona, output_dir=tmp_path)
        assert path.exists()
        assert path.stat().st_size > 0

    def test_handout_teacher_name_in_header(self, sample_lesson, sample_persona, tmp_path):
        from docx import Document

        from eduagent.doc_export import export_student_handout

        path = export_student_handout(sample_lesson, sample_persona, output_dir=tmp_path)
        doc = Document(str(path))
        full_text = "\n".join(p.text for p in doc.paragraphs)
        assert "Ms. Rivera" in full_text
