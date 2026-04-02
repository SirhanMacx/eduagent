"""PowerPoint (PPTX) export for lesson plans.

Generates professional, subject-themed slide decks with academic images.
"""

from __future__ import annotations

import asyncio
import logging
import re
from datetime import date
from pathlib import Path
from typing import TYPE_CHECKING, Optional

from clawed.async_utils import run_async_safe
from clawed.export_theme import _hex_to_rgb, _resolve_output

if TYPE_CHECKING:
    from clawed.models import DailyLesson, TeacherPersona

logger = logging.getLogger(__name__)

# Track temporary image files created during PPTX export for cleanup
_temp_image_files: list[Path] = []


def cleanup_temp_images() -> None:
    """Remove temporary image files created during PPTX export."""
    for tmp in _temp_image_files:
        try:
            tmp.unlink(missing_ok=True)
        except OSError:
            pass
    _temp_image_files.clear()


# ── PPTX helpers ──────────────────────────────────────────────────────


def _detect_subject(persona: "TeacherPersona") -> str:
    """Best-effort subject detection from persona fields."""
    subj = (persona.subject_area or "").strip().lower()
    if subj:
        return subj
    return ""


# ── Topic-based visual theming ──────────────────────────────────────

_TOPIC_THEMES: list[tuple[list[str], dict[str, str]]] = [
    # Exploration / nautical topics
    (
        ["exploration", "nautical", "voyage", "navigation", "maritime",
         "columbus", "magellan", "ocean", "sailing", "expedition", "discover"],
        {
            "primary": "1a3a5c",       # Navy blue
            "secondary": "c9a96e",     # Gold accent
            "accent": "e8eef5",        # Light blue-gray
            "bg_dark": "0f2440",       # Deep navy
            "bg_light": "f0f3f8",      # Soft blue-white
            "text_dark": "1A1A1A",
            "text_light": "FFFFFF",
            "text_dim": "888888",
        },
    ),
    # Renaissance / art topics
    (
        ["renaissance", "art", "painting", "sculpture", "michelangelo",
         "da vinci", "medici", "baroque", "classical art", "gallery",
         "fresco", "patron"],
        {
            "primary": "4a3728",       # Dark brown
            "secondary": "d4af37",     # Gold accent
            "accent": "f5efe0",        # Warm cream
            "bg_dark": "2e2018",       # Deep brown
            "bg_light": "faf6ef",      # Light cream
            "text_dark": "1A1A1A",
            "text_light": "FFFFFF",
            "text_dim": "888888",
        },
    ),
    # Revolution / war topics
    (
        ["revolution", "war", "battle", "conflict", "military", "independence",
         "civil war", "rebellion", "uprising", "army", "combat", "wwi", "wwii",
         "world war"],
        {
            "primary": "8b0000",       # Dark red
            "secondary": "ffd700",     # Gold accent
            "accent": "fce8e8",        # Light red tint
            "bg_dark": "4a0000",       # Deep red
            "bg_light": "fdf5f5",      # Soft rose
            "text_dark": "1A1A1A",
            "text_light": "FFFFFF",
            "text_dim": "888888",
        },
    ),
    # Science topics
    (
        ["science", "biology", "chemistry", "physics", "experiment",
         "hypothesis", "laboratory", "cell", "atom", "molecule",
         "ecosystem", "evolution", "genetics"],
        {
            "primary": "1a4a4a",       # Dark teal
            "secondary": "00c853",     # Bright green accent
            "accent": "e0f5ef",        # Light teal tint
            "bg_dark": "0f2e2e",       # Deep teal
            "bg_light": "f0faf5",      # Soft mint
            "text_dark": "1A1A1A",
            "text_light": "FFFFFF",
            "text_dim": "888888",
        },
    ),
    # Math topics
    (
        ["math", "algebra", "geometry", "calculus", "equation", "fraction",
         "polynomial", "trigonometry", "statistics", "probability",
         "arithmetic", "number"],
        {
            "primary": "2a1a4a",       # Dark purple
            "secondary": "ff9800",     # Orange accent
            "accent": "f0e8f8",        # Light purple tint
            "bg_dark": "1a0f30",       # Deep purple
            "bg_light": "f8f5fc",      # Soft lavender
            "text_dark": "1A1A1A",
            "text_light": "FFFFFF",
            "text_dim": "888888",
        },
    ),
]

_DEFAULT_TOPIC_THEME: dict[str, str] = {
    "primary": "2D5F3C",       # Deep green
    "secondary": "D4A843",     # Warm gold
    "accent": "EAF0E4",        # Soft green tint
    "bg_dark": "1A3D25",       # Dark green
    "bg_light": "F4F7F2",      # Light sage
    "text_dark": "1A1A1A",
    "text_light": "FFFFFF",
    "text_dim": "888888",
}


def get_topic_theme(title: str, subject: str = "") -> dict[str, str]:
    """Select a color theme based on lesson topic keywords.

    Scans the lesson title and subject for keywords associated with
    thematic palettes (nautical, renaissance, war, science, math).
    Returns the first matching theme or a default deep-green/gold palette.
    """
    combined = f"{title} {subject}".lower()
    for keywords, theme in _TOPIC_THEMES:
        for kw in keywords:
            if kw in combined:
                return theme
    return _DEFAULT_TOPIC_THEME


def _add_shape_fill(shape, hex_color: str) -> None:
    """Fill a shape with a solid color."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(hex_color)


def _set_text_props(run, font_size_pt: int, hex_color: str, bold: bool = False):
    """Set font properties on a text run."""
    from pptx.util import Pt

    run.font.size = Pt(font_size_pt)
    run.font.color.rgb = _hex_to_rgb(hex_color)
    run.font.bold = bold
    run.font.name = "Calibri"


def _split_text(text: str, max_len: int = 550) -> list[str]:
    """Split long text into chunks at sentence boundaries."""
    sentences = text.replace("\n", " ").split(". ")
    chunks: list[str] = []
    current = ""
    for s in sentences:
        candidate = f"{current}. {s}" if current else s
        if len(candidate) > max_len and current:
            chunks.append(current.strip())
            current = s
        else:
            current = candidate
    if current.strip():
        chunks.append(current.strip())
    return chunks or [text]


def _section_divider(prs, slide_num, text, theme, slide_w, slide_h,
                     sub_text: str = ""):
    """Create a section divider slide with label + optional sub-instruction."""
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Inches, Pt

    slide_num[0] += 1
    layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(layout)

    # Two-tone background: dark primary top 60%, light bottom 40%
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(theme.get("bg_dark", theme["primary"]))

    # Light bottom band
    band_top = int(slide_h * 0.62)
    band_h = slide_h - band_top
    band = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Emu(0), band_top, slide_w, band_h,
    )
    band.line.fill.background()
    bfill = band.fill
    bfill.solid()
    bfill.fore_color.rgb = _hex_to_rgb(theme.get("bg_light", "F0F0F0"))

    # Section label — left-aligned, 36pt, accent color, upper portion
    # Split "Let's Practice Together\n(15 minutes)" into label and time
    parts = text.split("\n", 1)
    label = parts[0].strip()
    time_hint = parts[1].strip() if len(parts) > 1 else ""

    tb = slide.shapes.add_textbox(
        Inches(1.0), Inches(1.8), slide_w - Inches(2.0), Inches(1.6),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = label
    run.font.size = Pt(38)
    run.font.color.rgb = _hex_to_rgb(theme.get("text_light", "FFFFFF"))
    run.font.bold = True
    run.font.name = "Calibri"

    # Time hint — smaller, below label
    if time_hint:
        tb_time = slide.shapes.add_textbox(
            Inches(1.0), Inches(3.5), Inches(4.0), Inches(0.6),
        )
        p_t = tb_time.text_frame.paragraphs[0]
        run_t = p_t.add_run()
        run_t.text = time_hint
        run_t.font.size = Pt(20)
        run_t.font.color.rgb = _hex_to_rgb(theme.get("text_light", "DDDDDD"))
        run_t.font.name = "Calibri"

    # Sub-instruction text (e.g. activity directions) in the light band
    if sub_text:
        tb_sub = slide.shapes.add_textbox(
            Inches(1.0), int(slide_h * 0.65), slide_w - Inches(2.0), int(slide_h * 0.30),
        )
        tf_sub = tb_sub.text_frame
        tf_sub.word_wrap = True
        p_sub = tf_sub.paragraphs[0]
        run_sub = p_sub.add_run()
        run_sub.text = sub_text
        run_sub.font.size = Pt(16)
        run_sub.font.color.rgb = _hex_to_rgb(theme.get("text_dark", "333333"))
        run_sub.font.name = "Calibri"

    # Footer with slide number
    left = slide_w - Inches(1.5)
    top = slide_h - Inches(0.45)
    tb_footer = slide.shapes.add_textbox(left, top, Inches(1.2), Inches(0.3))
    p_f = tb_footer.text_frame.paragraphs[0]
    p_f.alignment = PP_ALIGN.RIGHT
    run_f = p_f.add_run()
    run_f.text = str(slide_num[0])
    _set_text_props(run_f, 10, "999999")


# ── Image fetching ────────────────────────────────────────────────────


def _try_fetch_images(topics: list[tuple[str, str]], subject: str) -> dict[str, Optional[Path]]:
    """Attempt to fetch images for multiple topics. Non-blocking, short timeout.

    Returns a dict mapping key -> Path | None.
    """
    from clawed.slide_images import fetch_slide_image

    results: dict[str, Optional[Path]] = {}

    async def _fetch_all():
        for topic, key in topics:
            try:
                path = await asyncio.wait_for(
                    fetch_slide_image(topic, subject=subject),
                    timeout=5.0,
                )
                results[key] = path
            except Exception:
                results[key] = None

    try:
        run_async_safe(_fetch_all())
    except Exception as e:
        logger.debug("Image fetching failed: %s", e)

    return results


def _try_fetch_content_images(
    items: list[tuple[str, str, str]],
    subject: str,
    max_images: int = 4,
) -> dict[str, Optional[Path]]:
    """Fetch images based on slide *content* text, not just the lesson title.

    Each item is ``(content_text, fallback_topic, key)``.  Stops after
    ``max_images`` successful fetches to avoid slowing down generation.

    Returns a dict mapping key -> Path | None.
    """
    from clawed.slide_images import fetch_content_image, fetch_slide_image

    results: dict[str, Optional[Path]] = {}

    async def _fetch_all():
        found = 0
        for content_text, fallback_topic, key in items:
            if found >= max_images:
                results[key] = None
                continue
            try:
                if content_text:
                    path = await asyncio.wait_for(
                        fetch_content_image(
                            content_text,
                            subject=subject,
                            fallback_topic=fallback_topic,
                        ),
                        timeout=5.0,
                    )
                else:
                    # Title slide: use topic-based search
                    path = await asyncio.wait_for(
                        fetch_slide_image(fallback_topic, subject=subject),
                        timeout=5.0,
                    )
                results[key] = path
                if path:
                    found += 1
            except Exception:
                results[key] = None

    try:
        run_async_safe(_fetch_all())
    except Exception as e:
        logger.debug("Content image fetching failed: %s", e)

    return results


# ── Main export function ──────────────────────────────────────────────


def export_lesson_pptx(
    lesson: "DailyLesson",
    persona: "TeacherPersona",
    output_dir: Path | None = None,
    agent_name: str = "Claw-ED",
    include_images: bool = True,
) -> Path:
    """Generate a professional PowerPoint presentation from a lesson plan.

    Produces a polished, school-board-ready slide deck with:
    - Subject-themed color palette
    - Large readable fonts for projection on classroom screens
    - Academic images (LOC, Wikimedia, Unsplash) where available
    - Clean, modern, minimal design -- NOT busy or cluttered
    - Consistent visual language throughout

    Returns the path to the saved .pptx file.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Inches, Pt

    from clawed.sanitize import sanitize_text

    def _strip_inline_answers(text: str) -> str:
        """Remove (Answer: ...) and similar inline answer keys from student-facing slide text.

        The LLM sometimes embeds answer keys directly in guided practice /
        independent work fill-in-the-blank text.  These must never appear on
        projected slides — answers belong in speaker notes only.
        """
        if not text:
            return text
        # Remove (Answer: ...) and (answer: ...) patterns
        text = re.sub(r"\s*\((?:Answer|answer|Ans|ans)\s*:\s*[^)]+\)", "", text)
        # Remove [Answer: ...] bracket variant
        text = re.sub(r"\s*\[(?:Answer|answer|Ans|ans)\s*:\s*[^\]]+\]", "", text)
        # Remove inline answer keys like "_____ (Columbus)"  — parenthetical proper nouns after blanks
        # Only strip if the blank is present — don't strip all parentheticals
        text = re.sub(r"(_{2,})\s*\([A-Z][^)]{1,40}\)", r"\1", text)
        return text.strip()


    # Sanitize all lesson text fields before rendering to slides
    lesson.title = sanitize_text(lesson.title)
    lesson.objective = sanitize_text(lesson.objective)
    lesson.do_now = sanitize_text(lesson.do_now) if lesson.do_now else ""
    lesson.direct_instruction = sanitize_text(lesson.direct_instruction) if lesson.direct_instruction else ""
    # Strip answers before sanitizing — keep originals for speaker notes
    _gp_with_answers = lesson.guided_practice or ""
    _iw_with_answers = lesson.independent_work or ""
    lesson.guided_practice = _strip_inline_answers(sanitize_text(_gp_with_answers)) if _gp_with_answers else ""
    lesson.independent_work = _strip_inline_answers(sanitize_text(_iw_with_answers)) if _iw_with_answers else ""
    if lesson.homework:
        lesson.homework = sanitize_text(lesson.homework)
    for q in lesson.exit_ticket:
        q.question = sanitize_text(q.question)
        q.expected_response = sanitize_text(q.expected_response)
    lesson.standards = [sanitize_text(s) for s in lesson.standards]
    lesson.materials_needed = [sanitize_text(m) for m in lesson.materials_needed]

    def _clean_slide_text(text: str) -> str:
        """Final text cleanup before rendering to slides.

        Removes markdown artifacts, normalises blanks, strips section
        headers embedded in body text, and trims prose to slide-safe length.
        """
        if not text:
            return text
        # Replace __blank__ or ________ with a standard blank line
        text = re.sub(r"_{2,}", "______", text)
        # Strip markdown bold/italic markers
        text = re.sub(r"\*{1,3}([^*]+)\*{1,3}", r"\1", text)
        # Strip ## section headers at line start (DI prose has these)
        text = re.sub(r"^#{1,4}\s+", "", text, flags=re.MULTILINE)
        # Strip leading open-quote orphan chars (split source excerpts)
        text = re.sub(r'^["\u201c\u2018\u2019\u201d]\s*', "", text.strip())
        return text.strip()

    # Resolve teacher display name
    teacher_display_name = ""
    if persona and persona.name and persona.name != "My Teaching Persona":
        teacher_display_name = persona.name
    else:
        try:
            from clawed.models import AppConfig as _AppConfig
            _cfg = _AppConfig.load()
            if _cfg.teacher_profile and _cfg.teacher_profile.name:
                teacher_display_name = _cfg.teacher_profile.name
        except Exception:
            pass
    if not teacher_display_name:
        teacher_display_name = "Teacher"

    subject = _detect_subject(persona)
    # Topic-based visual theming: match lesson title/subject keywords to
    # curated palettes (nautical, renaissance, war, science, math).
    # Falls back to a default deep-green/gold theme.
    theme = get_topic_theme(lesson.title, subject)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # ── Try to fetch images asynchronously (entity-based queries) ────
    # Extracts named entities (people, places, documents) from lesson content
    # and searches for those specifically. Up to 5 images per deck.
    # Images on: Title (bg), Do Now (accent), Direct Instruction (sidebar)
    # No images on: Objectives, Guided Practice, Exit Ticket, Closing
    if include_images:
        from clawed.slide_images import _fetch_wikimedia, extract_image_subjects
        entities = extract_image_subjects(lesson)

        # For named entities (people, places, documents), go directly to Wikipedia —
        # it's faster and the article thumbnail is always the right image.
        # For generic topic queries, fall back to the standard LOC→Wikimedia chain.
        async def _fetch_entity_images() -> dict[str, Optional[Path]]:
            results: dict[str, Optional[Path]] = {}
            found = 0
            for i, entity in enumerate(entities[:5]):
                if found >= 5:
                    break
                key = f"entity_{i}"
                query = entity["query"]
                try:
                    path = await asyncio.wait_for(
                        _fetch_wikimedia(query),
                        timeout=8.0,
                    )
                    results[key] = path
                    if path:
                        found += 1
                        logger.info("Entity image [%s]: %s -> %s", key, query, path)
                    else:
                        logger.debug("No entity image for: %s", query)
                except Exception as e:
                    results[key] = None
                    logger.debug("Entity image fetch failed for %s: %s", query, e)
            return results

        if entities:
            try:
                images = run_async_safe(_fetch_entity_images())
            except Exception as e:
                logger.debug("Entity image fetch block failed: %s", e)
                images = {}
        else:
            # No entities found — use lesson title direct Wikipedia lookup
            try:
                async def _fetch_title_image() -> dict[str, Optional[Path]]:
                    path = await asyncio.wait_for(
                        _fetch_wikimedia(lesson.title),
                        timeout=8.0,
                    )
                    return {"entity_0": path}
                images = run_async_safe(_fetch_title_image())
            except Exception:
                images = {}

        logger.info(
            "Image fetch: %d entities, %d images found",
            len(entities),
            len([v for v in images.values() if v]),
        )
    else:
        images = {}
        logger.info("Generating PPTX with include_images=%s", include_images)

    # ── Shared layout helpers ─────────────────────────────────────────
    slide_num = [0]

    def _next_slide():
        slide_num[0] += 1
        layout = prs.slide_layouts[6]  # blank layout
        return prs.slides.add_slide(layout)

    def _add_footer(slide, num: int):
        left = slide_w - Inches(1.5)
        top = slide_h - Inches(0.45)
        tb = slide.shapes.add_textbox(left, top, Inches(1.2), Inches(0.3))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = str(num)
        _set_text_props(run, 10, "999999")

    def _bar(slide, left, top, width, height, hex_color: str):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height,
        )
        shape.line.fill.background()
        _add_shape_fill(shape, hex_color)
        return shape

    def _rounded_card(slide, left, top, width, height, hex_color: str):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
        )
        shape.line.fill.background()
        _add_shape_fill(shape, hex_color)
        # Subtle shadow
        try:
            shape.shadow.inherit = False
            shadow = shape.shadow
            shadow.visible = True
            shadow.blur_radius = Pt(4)
            shadow.distance = Pt(2)
        except Exception:
            pass
        return shape

    def _white_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb("FFFFFF")

    def _tinted_bg(slide, hex_color: str):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb(hex_color)

    def _safe_image_path(image_path: Path) -> Optional[Path]:
        """Validate an image file before embedding in a slide.

        python-pptx crashes on palette-mode PNGs, truncated files, or
        non-image data saved with an image extension.  This converts
        palette/RGBA images to RGB and rejects anything unreadable.

        Returns a safe path (possibly a converted temp file) or None.
        """
        try:
            import io
            import tempfile

            from PIL import Image as PILImage
            data = image_path.read_bytes()
            if len(data) < 1000:
                return None
            img = PILImage.open(io.BytesIO(data))
            # Convert palette (P) or RGBA to RGB — python-pptx can't handle these
            if img.mode in ("P", "RGBA", "LA", "PA"):
                img = img.convert("RGB")
                tmp = Path(tempfile.mktemp(suffix=".jpg"))
                img.save(tmp, "JPEG", quality=90)
                _temp_image_files.append(tmp)
                return tmp
            return image_path
        except Exception as e:
            logger.debug("Image validation failed for %s: %s", image_path, e)
            return None

    def _add_bg_image(slide, image_path: Path, overlay_alpha: str = "30000"):
        """Full-bleed background image with dark gradient overlay."""
        safe = _safe_image_path(image_path)
        if not safe:
            return
        try:
            pic = slide.shapes.add_picture(
                str(safe), Emu(0), Emu(0), slide_w, slide_h,
            )
        except Exception as e:
            logger.debug("_add_bg_image failed: %s", e)
            return
        sp = pic._element
        sp.getparent().remove(sp)
        slide.shapes._spTree.insert(2, sp)

        from pptx.oxml.ns import qn

        overlay = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), slide_w, slide_h,
        )
        overlay.line.fill.background()
        fill_obj = overlay.fill
        fill_obj.solid()
        fill_obj.fore_color.rgb = _hex_to_rgb(theme.get("bg_dark", theme["primary"]))
        solid_fill = overlay._element.spPr.solidFill
        if solid_fill is not None:
            srgb = solid_fill.find(qn("a:srgbClr"))
            if srgb is None:
                srgb = solid_fill.find(qn("a:sysClr"))
            if srgb is not None:
                from lxml import etree
                alpha = etree.SubElement(srgb, qn("a:alpha"))
                alpha.set("val", overlay_alpha)

    def _add_sidebar_image(slide, image_path: Path, caption: str = ""):
        safe = _safe_image_path(image_path)
        if not safe:
            return
        img_left = int(slide_w * 0.65)
        img_width = int(slide_w * 0.33)
        img_top = Inches(1.4)
        img_height = int(slide_h - Inches(2.4))
        try:
            slide.shapes.add_picture(
                str(safe), img_left, img_top, img_width, img_height,
            )
            if caption:
                cap_tb = slide.shapes.add_textbox(
                    img_left, img_top + img_height + Inches(0.05),
                    img_width, Inches(0.35),
                )
                cap_p = cap_tb.text_frame.paragraphs[0]
                cap_p.alignment = PP_ALIGN.CENTER
                cap_run = cap_p.add_run()
                cap_run.text = caption
                _set_text_props(cap_run, 9, "888888")
                cap_run.font.italic = True
        except Exception:
            pass

    def _add_accent_image(slide, image_path: Path, caption: str = ""):
        """Add a smaller accent image in the bottom-right area."""
        safe = _safe_image_path(image_path)
        if not safe:
            return
        img_left = slide_w - Inches(4.0)
        img_top = slide_h - Inches(3.2)
        img_width = Inches(3.5)
        img_height = Inches(2.2)
        try:
            slide.shapes.add_picture(
                str(safe), img_left, img_top, img_width, img_height,
            )
            if caption:
                cap_tb = slide.shapes.add_textbox(
                    img_left, img_top + img_height + Inches(0.05),
                    img_width, Inches(0.3),
                )
                cap_p = cap_tb.text_frame.paragraphs[0]
                cap_p.alignment = PP_ALIGN.CENTER
                cap_run = cap_p.add_run()
                cap_run.text = caption
                _set_text_props(cap_run, 9, "888888")
                cap_run.font.italic = True
        except Exception:
            pass

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 1: TITLE
    # Full-bleed image (or solid primary bg) with lesson title overlay
    # ═══════════════════════════════════════════════════════════════════
    slide = _next_slide()

    title_image = images.get("entity_0")
    if title_image:
        _add_bg_image(slide, title_image, overlay_alpha="30000")  # 70% dark overlay
    else:
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb(theme.get("bg_dark", theme["primary"]))
        # Gradient effect: lighter bar at bottom
        _bar(slide, Emu(0), int(slide_h * 0.7), slide_w, int(slide_h * 0.3), theme["primary"])

    # Lesson title -- 44pt white bold, centered
    tb = slide.shapes.add_textbox(
        Inches(1.5), Inches(1.5), slide_w - Inches(3.0), Inches(2.5),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = lesson.title
    _set_text_props(run, 44, theme["text_light"], bold=True)

    # Aim -- the driving question, centered below title
    aim_text = lesson.objective
    if aim_text and not aim_text.lower().startswith("students will"):
        aim_display = f"Aim: {aim_text}"
    elif aim_text:
        aim_display = aim_text
    else:
        aim_display = ""
    if aim_display:
        tb_aim = slide.shapes.add_textbox(
            Inches(1.5), Inches(4.0), slide_w - Inches(3.0), Inches(1.0),
        )
        tf_aim = tb_aim.text_frame
        tf_aim.word_wrap = True
        p_aim = tf_aim.paragraphs[0]
        p_aim.alignment = PP_ALIGN.CENTER
        run_aim = p_aim.add_run()
        run_aim.text = aim_display
        _set_text_props(run_aim, 20, "EEEEEE")
        run_aim.font.italic = True

    # Subtitle: "Teacher Name | Date" -- 16pt
    tb2 = slide.shapes.add_textbox(
        Inches(1.5), Inches(5.2), slide_w - Inches(3.0), Inches(0.8),
    )
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = (
        f"{teacher_display_name}  |  "
        f"{date.today().strftime('%B %d, %Y')}"
    )
    _set_text_props(run2, 16, "BBBBBB")

    # Bottom accent bar in subject color
    _bar(slide, Emu(0), slide_h - Inches(0.15), slide_w, Inches(0.15), theme["secondary"])

    _add_footer(slide, slide_num[0])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 2: OBJECTIVES
    # Clean white bg, left accent bar, SWBAT in colored card
    # ═══════════════════════════════════════════════════════════════════
    slide = _next_slide()
    _white_bg(slide)

    # Left accent bar (4px, subject color)
    _bar(slide, Inches(0.5), Inches(0.8), Inches(0.06), Inches(5.5), theme["primary"])

    # "Today's Objectives" header -- 36pt, subject color
    tb = slide.shapes.add_textbox(
        Inches(1.0), Inches(0.8), slide_w - Inches(2.0), Inches(1.0),
    )
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Today's Objectives"
    _set_text_props(run, 36, theme["primary"], bold=True)

    # SWBAT card -- colored rounded rectangle
    _rounded_card(
        slide,
        Inches(1.0), Inches(2.2),
        slide_w - Inches(2.0), Inches(1.8),
        theme["accent"],
    )
    tb = slide.shapes.add_textbox(
        Inches(1.4), Inches(2.5), slide_w - Inches(2.8), Inches(1.2),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.space_after = Pt(12)
    run = p.add_run()
    run.text = f"SWBAT: {lesson.objective}"
    _set_text_props(run, 24, theme["text_dark"], bold=True)

    # Standards below in 16pt gray
    if lesson.standards:
        tb = slide.shapes.add_textbox(
            Inches(1.0), Inches(4.4), slide_w - Inches(2.0), Inches(2.5),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Standards Addressed"
        _set_text_props(run, 16, "888888", bold=True)

        for std in lesson.standards[:3]:
            p = tf.add_paragraph()
            p.space_before = Pt(6)
            run = p.add_run()
            run.text = f"  {std}"
            _set_text_props(run, 16, "666666")

    _add_footer(slide, slide_num[0])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 3: DO NOW / WARM-UP
    # Light tint bg, "Do Now" badge, question text, timer
    # ═══════════════════════════════════════════════════════════════════
    if lesson.do_now:
        slide = _next_slide()
        _tinted_bg(slide, theme["accent"])

        # "Do Now" badge -- rounded rectangle, subject color, white text
        badge = _rounded_card(
            slide,
            Inches(0.8), Inches(0.6),
            Inches(2.5), Inches(0.7),
            theme["primary"],
        )
        badge_tf = badge.text_frame
        badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = badge_tf.paragraphs[0].add_run()
        run.text = "Do Now"
        _set_text_props(run, 22, theme["text_light"], bold=True)

        # Question / prompt text — adaptive font size based on length
        # Do Now uses accent image (small, bottom-right) not sidebar
        do_now_img = images.get("entity_1")

        tb = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.8), slide_w - Inches(2.0), Inches(4.2),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.line_spacing = Pt(30)
        run = p.add_run()
        # Brief prompt on slide face, full text in speaker notes
        dn_text = _clean_slide_text(lesson.do_now)
        if len(dn_text) > 300:
            cutoff = dn_text[:300].rfind(". ")
            dn_display = dn_text[:cutoff + 1] if cutoff > 80 else dn_text[:300].rsplit(" ", 1)[0] + "..."
        else:
            dn_display = dn_text
        run.text = dn_display
        # Scale font: short prompt = 22pt, medium = 18pt, long = 16pt
        dn_font = 22 if len(dn_display) < 120 else 18 if len(dn_display) < 250 else 16
        _set_text_props(run, dn_font, theme["text_dark"])

        # Full Do Now in speaker notes
        if len(dn_text) > 250:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = dn_text

        if do_now_img:
            from clawed.slide_images import _extract_key_concepts
            concepts = _extract_key_concepts(lesson.do_now)
            caption = ", ".join(concepts[:2]) if concepts else ""
            _add_accent_image(slide, do_now_img, caption=caption)

        # Timer indicator
        minutes = lesson.time_estimates.get("do_now", 5)
        tb_timer = slide.shapes.add_textbox(
            slide_w - Inches(2.5), slide_h - Inches(1.0),
            Inches(2.0), Inches(0.5),
        )
        p_t = tb_timer.text_frame.paragraphs[0]
        p_t.alignment = PP_ALIGN.RIGHT
        run_t = p_t.add_run()
        run_t.text = f"{minutes} minutes"
        _set_text_props(run_t, 16, "888888")

        _add_footer(slide, slide_num[0])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 4: DIRECT INSTRUCTION (student-facing summary)
    # Full script goes into speaker notes; slide face shows brief summary.
    # Additional slides for vocabulary terms and source excerpts if found.
    # ═══════════════════════════════════════════════════════════════════
    if lesson.direct_instruction:
        di_text = lesson.direct_instruction

        # ── Summary slide ──────────────────────────────────────────
        slide = _next_slide()
        _white_bg(slide)

        # "Direct Instruction" badge (subject color, rounded)
        badge = _rounded_card(
            slide,
            Inches(0.8), Inches(0.6),
            Inches(4.5), Inches(0.7),
            theme["primary"],
        )
        badge_tf = badge.text_frame
        badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = badge_tf.paragraphs[0].add_run()
        run.text = "Direct Instruction"
        _set_text_props(run, 22, theme["text_light"], bold=True)

        # Side accent bar
        _bar(slide, Inches(0.6), Inches(1.7), Inches(0.06), Inches(5.0), theme["secondary"])

        # Sidebar image if available
        img_path = images.get("entity_2")
        text_width = slide_w - Inches(2.0)
        if img_path:
            text_width = int(slide_w * 0.60)
            from clawed.slide_images import _extract_key_concepts
            concepts = _extract_key_concepts(di_text)
            caption = ", ".join(concepts[:2]) if concepts else ""
            _add_sidebar_image(slide, img_path, caption=caption)

        # Extract 3-4 key bullet points from DI prose for slide face.
        # Full script goes to speaker notes.
        di_clean = _clean_slide_text(di_text)

        def _extract_di_bullets(text: str, max_bullets: int = 4) -> list[str]:
            """Pull key sentences/bullets from direct instruction prose."""
            bullets: list[str] = []
            # Try existing bullet/numbered list items first
            list_items = re.findall(r"^[\-•\*]\s+(.+)$", text, re.MULTILINE)
            if not list_items:
                list_items = re.findall(r"^\d+[\.\)]\s+(.+)$", text, re.MULTILINE)
            for item in list_items[:max_bullets]:
                item = item.strip()
                if len(item) > 20:
                    # Truncate at word boundary, not mid-word
                    if len(item) > 140:
                        item = item[:140].rsplit(' ', 1)[0] + '...'
                    bullets.append(item)
            if bullets:
                return bullets
            # Fallback: first N sentences that are short enough to be readable
            sentences = re.split(r"(?<=[.!?])\s+", text)
            for sent in sentences:
                sent = sent.strip()
                # Skip header-like short sentences and very long ones
                if 30 < len(sent) <= 160 and not sent.startswith(("##", "**", "Alright")):
                    bullets.append(sent)
                if len(bullets) >= max_bullets:
                    break
            return bullets or [text[:160].rsplit(" ", 1)[0] + "…"]

        bullets = _extract_di_bullets(di_clean)

        tb = slide.shapes.add_textbox(
            Inches(1.0), Inches(1.8), text_width, Inches(4.8),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        for bi, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if bi == 0 else tf.add_paragraph()
            p.line_spacing = Pt(26)
            p.space_before = Pt(6)
            run = p.add_run()
            run.text = f"• {bullet}"
            _set_text_props(run, 18, theme["text_dark"])

        # Full DI script in speaker notes
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = di_text

        _add_footer(slide, slide_num[0])

        # ── Vocabulary slide ───────────────────────────────────────
        # Prefer structured vocabulary from v2.3 model fields
        vocab_pairs: list[tuple[str, str]] = []
        if hasattr(lesson, "vocabulary") and lesson.vocabulary:
            for vt in lesson.vocabulary:
                term = getattr(vt, "term", "") if hasattr(vt, "term") else vt.get("term", "")
                defn = getattr(vt, "definition", "") if hasattr(vt, "definition") else vt.get("definition", "")
                if term and defn:
                    vocab_pairs.append((term, defn))

        # Fallback: regex extraction from DI text
        if not vocab_pairs:
            di_text_raw = lesson.direct_instruction or ""
            vocab_patterns = re.findall(
                r'\*\*([A-Z][^*]{2,30})\*\*\s*[-—:]+\s*([^*\n]{15,200})',
                di_text_raw,
            )
            if not vocab_patterns:
                vocab_patterns = re.findall(
                    r'(?:^|\n)\s*\*?\*?([A-Z][a-z]+(?:\s+[A-Za-z]+){0,3})\*?\*?\s*[-—:]+\s+'
                    r'([A-Za-z][^\n]{15,200})',
                    di_text_raw,
                )
            INSTRUCTIONAL_WORDS = {
                "check", "ask", "call", "facilitate", "transition", "minutes",
                "discuss", "display", "distribute", "brief", "move", "moved",
                "students", "responses", "excellent", "turn", "now", "let",
                "today", "good", "morning", "scholars", "friends", "class",
                "next", "first", "take", "look", "read", "write", "complete",
            }
            vocab_pairs = [
                (term.strip(), defn.strip()) for term, defn in vocab_patterns
                if 1 <= len(term.split()) <= 4
                and len(defn.split()) >= 3
                and not any(w in defn.lower().split()[:3] for w in INSTRUCTIONAL_WORDS)
                and not any(w in term.lower().split() for w in INSTRUCTIONAL_WORDS)
            ]

        if vocab_pairs:
            # Split vocabulary across multiple slides (max 4 terms per slide)
            # to prevent overcrowding when lessons have many terms.
            VOCAB_PER_SLIDE = 4
            vocab_chunks = [
                vocab_pairs[i:i + VOCAB_PER_SLIDE]
                for i in range(0, len(vocab_pairs), VOCAB_PER_SLIDE)
            ]
            total_vocab_slides = len(vocab_chunks)

            for chunk_idx, chunk in enumerate(vocab_chunks):
                slide = _next_slide()
                _white_bg(slide)

                # "Key Vocabulary" badge (with page indicator if multi-slide)
                badge_label = "Key Vocabulary"
                if total_vocab_slides > 1:
                    badge_label = f"Key Vocabulary ({chunk_idx + 1}/{total_vocab_slides})"
                badge = _rounded_card(
                    slide,
                    Inches(0.8), Inches(0.6),
                    Inches(3.5) if total_vocab_slides == 1 else Inches(4.5),
                    Inches(0.7),
                    theme["secondary"],
                )
                badge_tf = badge.text_frame
                badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                run = badge_tf.paragraphs[0].add_run()
                run.text = badge_label
                _set_text_props(run, 22, theme["text_light"], bold=True)

                # Left accent bar
                _bar(slide, Inches(0.6), Inches(1.7), Inches(0.06), Inches(5.0), theme["primary"])

                # Vocabulary: 2-column layout — term (bold, left) | definition (right)
                LEFT_COL_W = Inches(3.8)
                RIGHT_COL_W = slide_w - Inches(5.2)
                # With max 4 terms, use more vertical space per row
                row_h = Inches(1.1)
                start_y = Inches(1.75)

                for idx, (term, definition) in enumerate(chunk):
                    y = start_y + idx * row_h
                    if y + row_h > slide_h - Inches(0.5):
                        break  # safety: don't overflow bottom

                    # Term — left column, bold, primary color
                    tb_term = slide.shapes.add_textbox(
                        Inches(0.9), y, LEFT_COL_W, row_h,
                    )
                    tf_term = tb_term.text_frame
                    tf_term.word_wrap = True
                    p_term = tf_term.paragraphs[0]
                    p_term.line_spacing = Pt(22)
                    run_t = p_term.add_run()
                    run_t.text = term.strip()
                    _set_text_props(run_t, 18, theme["primary"], bold=True)

                    # Definition — right column, normal weight
                    tb_def = slide.shapes.add_textbox(
                        Inches(4.9), y, RIGHT_COL_W, row_h,
                    )
                    tf_def = tb_def.text_frame
                    tf_def.word_wrap = True
                    p_def = tf_def.paragraphs[0]
                    p_def.line_spacing = Pt(22)
                    run_d = p_def.add_run()
                    # Truncate long definitions to keep them slide-readable
                    defn = definition.strip()
                    if len(defn) > 120:
                        defn = defn[:117].rsplit(" ", 1)[0] + "…"
                    run_d.text = defn
                    _set_text_props(run_d, 18, theme["text_dark"])

                    # Light separator line between rows
                    if idx < len(chunk) - 1:
                        _bar(slide, Inches(0.9), y + row_h - Inches(0.04),
                             slide_w - Inches(1.8), Inches(0.02), "EEEEEE")

                _add_footer(slide, slide_num[0])

        # ── Primary Source slides (from structured lesson data) ─────
        # Use lesson.primary_sources (PrimarySourceDocument objects)
        # instead of regex-extracting fragments from DI prose.
        source_img_idx = 3
        primary_sources = getattr(lesson, "primary_sources", []) or []
        for ps in primary_sources[:3]:
            # Get source fields — handle both dict and object
            if isinstance(ps, dict):
                ps_title = ps.get("title", "") or ps.get("document_label", "")
                ps_author = ps.get("author", "")
                ps_date = ps.get("date", "")
                ps_context = ps.get("context", "")
                ps_text = ps.get("full_text", "") or ps.get("text", "")
                ps_questions = ps.get("analysis_questions", [])
            else:
                ps_title = getattr(ps, "title", "") or getattr(ps, "document_label", "")
                ps_author = getattr(ps, "author", "")
                ps_date = getattr(ps, "date", "")
                ps_context = getattr(ps, "context", "")  # noqa: F841
                ps_text = getattr(ps, "full_text", "")
                ps_questions = getattr(ps, "analysis_questions", [])

            if not ps_text and not ps_title:
                continue

            slide = _next_slide()
            _tinted_bg(slide, theme["accent"])

            # "Primary Source" badge with title
            badge_text = ps_title or "Primary Source"
            badge = _rounded_card(
                slide,
                Inches(0.8), Inches(0.5),
                Inches(5.0), Inches(0.7),
                theme["primary"],
            )
            badge_tf = badge.text_frame
            badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = badge_tf.paragraphs[0].add_run()
            run.text = _clean_slide_text(badge_text)[:80]
            _set_text_props(run, 20, theme["text_light"], bold=True)

            # Attribution line (author + date)
            attribution = ""
            if ps_author:
                attribution = ps_author
            if ps_date:
                attribution += f", {ps_date}" if attribution else ps_date
            if attribution:
                tb_attr = slide.shapes.add_textbox(
                    Inches(0.8), Inches(1.3), Inches(8.0), Inches(0.4),
                )
                p_attr = tb_attr.text_frame.paragraphs[0]
                run_attr = p_attr.add_run()
                run_attr.text = f"— {_clean_slide_text(attribution)}"
                _set_text_props(run_attr, 14, theme["text_dim"])
                run_attr.font.italic = True

            # Source text — the actual excerpt
            src_img = images.get(f"entity_{source_img_idx}")
            text_width = slide_w - Inches(2.0)
            if src_img:
                text_width = int(slide_w * 0.58)
                _add_sidebar_image(slide, src_img)
            source_img_idx += 1

            clean_text = _clean_slide_text(ps_text.strip())
            if len(clean_text) > 400:
                clean_text = clean_text[:397].rsplit(" ", 1)[0] + "…"

            text_font = 18 if len(clean_text) < 200 else 15 if len(clean_text) < 350 else 13

            y_start = Inches(1.8) if attribution else Inches(1.5)
            tb = slide.shapes.add_textbox(
                Inches(1.0), y_start, text_width, Inches(3.5),
            )
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = Pt(text_font + 7)
            run = p.add_run()
            run.text = f"\u201c{clean_text}\u201d" if clean_text else "(Source text not available)"
            _set_text_props(run, text_font, theme["text_dark"])
            run.font.italic = True

            # Analysis questions (if any)
            if ps_questions:
                q_y = slide_h - Inches(2.0)
                tb_q = slide.shapes.add_textbox(
                    Inches(1.0), q_y, slide_w - Inches(2.0), Inches(1.2),
                )
                tf_q = tb_q.text_frame
                tf_q.word_wrap = True
                for qi, question in enumerate(ps_questions[:2]):
                    p_q = tf_q.paragraphs[0] if qi == 0 else tf_q.add_paragraph()
                    p_q.space_before = Pt(4)
                    run_q = p_q.add_run()
                    q_text = _clean_slide_text(question)
                    if len(q_text) > 100:
                        q_text = q_text[:97].rsplit(" ", 1)[0] + "…"
                    run_q.text = f"Q{qi+1}: {q_text}"
                    _set_text_props(run_q, 14, theme["text_dark"], bold=True)

            _add_footer(slide, slide_num[0])

    # ═══════════════════════════════════════════════════════════════════
    # SECTION DIVIDER: "Let's Practice Together"
    # ═══════════════════════════════════════════════════════════════════
    if lesson.guided_practice:
        gp_min = lesson.time_estimates.get("guided_practice", 15)
        # Show first instruction line in the divider's light band
        gp_preview = _clean_slide_text(lesson.guided_practice)
        gp_first_line = gp_preview.split("\n")[0].strip()
        if len(gp_first_line) > 120:
            gp_first_line = gp_first_line[:117].rsplit(" ", 1)[0] + "…"
        _section_divider(
            prs, slide_num, f"Let's Practice Together\n({gp_min} minutes)",
            theme, slide_w, slide_h,
            sub_text=gp_first_line,
        )

    # ═══════════════════════════════════════════════════════════════════
    # GUIDED PRACTICE -- "Your Turn" header
    # ═══════════════════════════════════════════════════════════════════
    if lesson.guided_practice:
        slide = _next_slide()
        _tinted_bg(slide, "F5F5F5")  # very light gray

        # "Your Turn" badge
        badge = _rounded_card(
            slide,
            Inches(0.8), Inches(0.6),
            Inches(3.0), Inches(0.7),
            theme["secondary"],
        )
        badge_tf = badge.text_frame
        badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = badge_tf.paragraphs[0].add_run()
        run.text = "Your Turn"
        _set_text_props(run, 22, theme["text_light"], bold=True)

        # No image on Guided Practice -- keep clean for readability

        # Guided practice: render as numbered/bulleted list items at 18pt
        gp_text = _clean_slide_text(lesson.guided_practice)

        # Split into individual items (lines starting with - or numbers)
        gp_lines = [line.strip() for line in gp_text.split("\n") if line.strip()]
        # Limit to what fits: ~7 lines at 18pt in 4.5" box
        MAX_GP_LINES = 7
        if len(gp_lines) > MAX_GP_LINES:
            gp_lines = gp_lines[:MAX_GP_LINES]

        tb = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.8), slide_w - Inches(2.0), Inches(4.8),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        for li, line in enumerate(gp_lines):
            p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            p.line_spacing = Pt(26)
            p.space_before = Pt(4)
            run = p.add_run()
            # Ensure list items have a bullet if they don't already
            if not line.startswith(("-", "•", "–")) and not re.match(r"^\d+[\.\)]", line):
                line = f"• {line}"
            run.text = line
            _set_text_props(run, 18, theme["text_dark"])

        # Full guided practice + ANSWER KEY in speaker notes
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = "ANSWER KEY:\n" + _gp_with_answers + "\n\n---\nSTUDENT VIEW:\n" + gp_text

        # Time estimate
        minutes = lesson.time_estimates.get("guided_practice", 15)
        tb_time = slide.shapes.add_textbox(
            slide_w - Inches(2.5), slide_h - Inches(1.0),
            Inches(2.0), Inches(0.5),
        )
        p_t = tb_time.text_frame.paragraphs[0]
        p_t.alignment = PP_ALIGN.RIGHT
        run_t = p_t.add_run()
        run_t.text = f"{minutes} minutes"
        _set_text_props(run_t, 16, "888888")

        _add_footer(slide, slide_num[0])

    # ═══════════════════════════════════════════════════════════════════
    # INDEPENDENT WORK -- "Independent Practice" header
    # ═══════════════════════════════════════════════════════════════════
    if lesson.independent_work:
        slide = _next_slide()
        _white_bg(slide)

        # "Independent Practice" badge
        badge = _rounded_card(
            slide,
            Inches(0.8), Inches(0.6),
            Inches(4.5), Inches(0.7),
            theme["primary"],
        )
        badge_tf = badge.text_frame
        badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = badge_tf.paragraphs[0].add_run()
        run.text = "Independent Practice"
        _set_text_props(run, 22, theme["text_light"], bold=True)

        # Left accent bar
        _bar(slide, Inches(0.6), Inches(1.7), Inches(0.06), Inches(5.0), theme["secondary"])

        # No image on Independent Work -- keep clean for readability

        # Brief student-facing instructions on slide face
        iw_text = lesson.independent_work
        iw_summary = iw_text
        if len(iw_text) > 250:
            cutoff = iw_text[:250].rfind(". ")
            if cutoff > 80:
                iw_summary = iw_text[:cutoff + 1]
            else:
                iw_summary = iw_text[:250].rsplit(" ", 1)[0] + "..."

        tb = slide.shapes.add_textbox(
            Inches(1.0), Inches(1.8), slide_w - Inches(2.0), Inches(4.5),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.line_spacing = Pt(34)
        run = p.add_run()
        run.text = iw_summary
        _set_text_props(run, 24, theme["text_dark"])

        # Full independent work in speaker notes
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = iw_text

        # Time estimate in corner
        minutes = lesson.time_estimates.get("independent_work", 10)
        _rounded_card(
            slide,
            slide_w - Inches(3.0), slide_h - Inches(1.2),
            Inches(2.2), Inches(0.6),
            theme["accent"],
        )
        tb_time = slide.shapes.add_textbox(
            slide_w - Inches(3.0), slide_h - Inches(1.2),
            Inches(2.2), Inches(0.6),
        )
        p_t = tb_time.text_frame.paragraphs[0]
        p_t.alignment = PP_ALIGN.CENTER
        run_t = p_t.add_run()
        run_t.text = f"{minutes} minutes"
        _set_text_props(run_t, 16, theme["text_dark"])

        _add_footer(slide, slide_num[0])

    # ═══════════════════════════════════════════════════════════════════
    # SECTION DIVIDER: "Show What You Know"
    # ═══════════════════════════════════════════════════════════════════
    if lesson.exit_ticket:
        _section_divider(prs, slide_num, "Show What You Know\n(Exit Ticket)", theme, slide_w, slide_h)

    # ═══════════════════════════════════════════════════════════════════
    # EXIT TICKET -- distinctive assessment design
    # Colored banner header, numbered question cards with shadow
    # ═══════════════════════════════════════════════════════════════════
    if lesson.exit_ticket:
        slide = _next_slide()
        _white_bg(slide)

        # Full-width colored banner for "Exit Ticket"
        _bar(slide, Emu(0), Emu(0), slide_w, Inches(1.2), theme["primary"])
        tb = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.15), slide_w - Inches(1.6), Inches(0.9),
        )
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "Exit Ticket"
        _set_text_props(run, 36, theme["text_light"], bold=True)

        # Numbered questions + answer line
        # Available height: 7.5" - 1.2" header - 0.6" footer = ~5.7"
        # Divide evenly among questions (max 4 shown)
        et_questions = lesson.exit_ticket[:4]
        n_q = len(et_questions)
        available_h = slide_h - Inches(1.55) - Inches(0.8)
        card_h = available_h // n_q

        q_top = Inches(1.5)
        for i, q in enumerate(et_questions, 1):
            # Card background
            _rounded_card(
                slide,
                Inches(0.8), q_top,
                slide_w - Inches(1.6), card_h - Inches(0.1),
                theme["accent"],
            )

            # Number circle — properly sized
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(1.1), q_top + Inches(0.12),
                Inches(0.55), Inches(0.55),
            )
            circle.line.fill.background()
            _add_shape_fill(circle, theme["secondary"])
            circ_p = circle.text_frame.paragraphs[0]
            circ_p.alignment = PP_ALIGN.CENTER
            run = circ_p.add_run()
            run.text = str(i)
            _set_text_props(run, 18, theme["text_light"], bold=True)

            # Question text — adaptive font, left of number
            q_text = _clean_slide_text(q.question)
            # For slide: trim to one clear sentence
            if len(q_text) > 150:
                q_text = q_text[:150].rsplit(" ", 1)[0] + "…"
            q_font = 16 if len(q_text) > 100 else 18

            tb = slide.shapes.add_textbox(
                Inches(1.85), q_top + Inches(0.10),
                slide_w - Inches(3.0), Inches(0.65),
            )
            tf = tb.text_frame
            tf.word_wrap = True
            run = tf.paragraphs[0].add_run()
            run.text = q_text
            _set_text_props(run, q_font, theme["text_dark"])

            # Answer line underneath question
            answer_y = q_top + Inches(0.75)
            _bar(slide, Inches(1.85), answer_y,
                 slide_w - Inches(3.0), Inches(0.02), "CCCCCC")

            q_top += card_h

        # "Turn in before you leave" footer note
        tb = slide.shapes.add_textbox(
            Inches(0.8), slide_h - Inches(0.75),
            slide_w - Inches(1.6), Inches(0.4),
        )
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "Turn in before you leave  •  REMEMBER: T.E.A."
        _set_text_props(run, 13, "888888")

        _add_footer(slide, slide_num[0])

    # ═══════════════════════════════════════════════════════════════════
    # CLOSING SLIDE -- Homework or Key Takeaway
    # ═══════════════════════════════════════════════════════════════════
    slide = _next_slide()
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(theme.get("bg_dark", theme["primary"]))

    if lesson.homework:
        # "Tonight's Homework" with details
        tb = slide.shapes.add_textbox(
            Inches(1.5), Inches(1.5), slide_w - Inches(3.0), Inches(1.5),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "Tonight's Homework"
        _set_text_props(run, 36, theme["text_light"], bold=True)

        _bar(
            slide, Inches(1.5), Inches(3.2),
            Inches(4.0), Inches(0.06), theme["secondary"],
        )

        tb = slide.shapes.add_textbox(
            Inches(1.5), Inches(3.6), slide_w - Inches(3.0), Inches(3.0),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.line_spacing = Pt(30)
        run = p.add_run()
        # Brief homework summary on slide, full in notes
        hw_text = lesson.homework
        if len(hw_text) > 200:
            cutoff = hw_text[:200].rfind(". ")
            hw_display = hw_text[:cutoff + 1] if cutoff > 60 else hw_text[:200].rsplit(" ", 1)[0] + "..."
        else:
            hw_display = hw_text
        run.text = hw_display
        _set_text_props(run, 22, "DDDDDD")
        if len(hw_text) > 200:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = f"HOMEWORK (full text):\n{hw_text}"
    else:
        # "Key Takeaway" or "Questions?"
        tb = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.5), slide_w - Inches(3.0), Inches(2.0),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "Questions?"
        _set_text_props(run, 48, theme["text_light"], bold=True)

    # Teacher name and subject in footer
    teacher_subject = subject.title() if subject else "Education"
    tb = slide.shapes.add_textbox(
        Inches(1.0), slide_h - Inches(1.0),
        slide_w - Inches(2.0), Inches(0.5),
    )
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"{teacher_display_name}  |  {teacher_subject}"
    _set_text_props(run, 14, "888888")

    # Watermark
    tb_wm = slide.shapes.add_textbox(
        Inches(1.0), slide_h - Inches(0.55),
        slide_w - Inches(2.0), Inches(0.3),
    )
    p_wm = tb_wm.text_frame.paragraphs[0]
    p_wm.alignment = PP_ALIGN.CENTER
    run_wm = p_wm.add_run()
    run_wm.text = f"Generated by {agent_name}"
    _set_text_props(run_wm, 10, "666666")

    _add_footer(slide, slide_num[0])

    # ── Save ──────────────────────────────────────────────────────────
    out = _resolve_output(output_dir, lesson, ".pptx")
    prs.save(str(out))
    cleanup_temp_images()
    return out
