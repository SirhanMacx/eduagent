"""Post-generation output review agent.

Checks every generated output (PPTX, DOCX, HTML game) for quality
issues before delivery to the teacher. Catches:
- Truncated/nonsensical text
- Images not matching slide content
- Answer keys visible in student materials
- Overcrowded slides
- Empty or incomplete sections
- Text starting mid-sentence

If issues are found, returns a structured report. The caller decides
whether to regenerate or deliver with warnings.
"""

from __future__ import annotations

import logging
import re
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)


class OutputReview:
    """Quality review results for a generated output."""

    def __init__(self) -> None:
        self.issues: list[dict[str, Any]] = []
        self.score: float = 10.0  # Start perfect, deduct for issues

    def add_issue(
        self, severity: str, location: str, description: str
    ) -> None:
        """Add a quality issue. Severity: critical, major, minor."""
        deductions = {"critical": 3.0, "major": 1.5, "minor": 0.5}
        self.issues.append({
            "severity": severity,
            "location": location,
            "description": description,
        })
        self.score = max(0.0, self.score - deductions.get(severity, 0.5))

    @property
    def passed(self) -> bool:
        """Output passes review if score >= 6.0 and no critical issues."""
        has_critical = any(i["severity"] == "critical" for i in self.issues)
        return self.score >= 6.0 and not has_critical

    def summary(self) -> str:
        """Human-readable summary."""
        if not self.issues:
            return f"PASSED ({self.score:.1f}/10) — no issues found"
        lines = [f"{'PASSED' if self.passed else 'FAILED'} ({self.score:.1f}/10)"]
        for issue in self.issues:
            icon = {"critical": "X", "major": "!", "minor": "~"}
            lines.append(
                f"  [{icon.get(issue['severity'], '?')}] "
                f"{issue['location']}: {issue['description']}"
            )
        return "\n".join(lines)


def review_pptx(pptx_path: Path) -> OutputReview:
    """Review a generated PPTX for quality issues."""
    from pptx import Presentation

    review = OutputReview()

    try:
        prs = Presentation(str(pptx_path))
    except Exception as e:
        review.add_issue("critical", "file", f"Cannot open PPTX: {e}")
        return review

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        loc = f"Slide {slide_num}"

        # Collect all text from the slide
        texts = []
        has_image = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
            if shape.shape_type == 13:
                has_image = True

        all_text = " ".join(texts)

        # Check: empty slide (no text at all)
        if not texts or (len(texts) == 1 and len(texts[0]) < 5):
            review.add_issue("major", loc, "Slide is empty or has minimal content")

        # Check: text starts mid-sentence (fragment)
        for t in texts:
            if t and len(t) > 10:
                first_char = t[0]
                if first_char.islower() or first_char in ("'", "\u2019"):
                    review.add_issue(
                        "major", loc,
                        f"Text starts mid-sentence: '{t[:50]}...'"
                    )
                    break

        # Check: truncated text (ends with ... or ellipsis mid-word)
        for t in texts:
            if t.endswith("...") or t.endswith("\u2026"):
                # Check if truncation is mid-word
                before_ellipsis = t.rstrip(".\u2026 ")
                if before_ellipsis and before_ellipsis[-1].isalpha():
                    review.add_issue(
                        "minor", loc,
                        f"Text truncated mid-word: '...{t[-40:]}'"
                    )

        # Check: too much text on one slide (>500 chars)
        if len(all_text) > 600:
            review.add_issue(
                "major", loc,
                f"Overcrowded slide ({len(all_text)} chars) — "
                f"should be split across multiple slides"
            )

        # Check: answer keys visible
        answer_patterns = [
            r"\(Answer:\s*[^)]+\)",
            r"\(answer:\s*[^)]+\)",
            r"ANSWER KEY",
            r"Answer:\s+\w+",
        ]
        for pattern in answer_patterns:
            if re.search(pattern, all_text):
                review.add_issue(
                    "critical", loc,
                    "Answer key visible in student-facing slide"
                )
                break

    return review


def review_docx(docx_path: Path) -> OutputReview:
    """Review a generated DOCX handout for quality issues."""
    from docx import Document

    review = OutputReview()

    try:
        doc = Document(str(docx_path))
    except Exception as e:
        review.add_issue("critical", "file", f"Cannot open DOCX: {e}")
        return review

    all_text = ""
    for para in doc.paragraphs:
        text = para.text.strip()
        all_text += text + " "

        # Check: answer keys visible
        if re.search(
            r"\(Answer:\s*[^)]+\)|\(answer:\s*[^)]+\)|ANSWER KEY",
            text,
        ):
            review.add_issue(
                "critical",
                "content",
                f"Answer key visible: '{text[:60]}...'"
            )

    # Check: document too short
    if len(all_text) < 500:
        review.add_issue(
            "major", "content", "Document is very short — may be incomplete"
        )

    # Check: missing key sections
    sections_expected = ["do now", "aim", "exit ticket"]
    text_lower = all_text.lower()
    for section in sections_expected:
        if section not in text_lower:
            review.add_issue(
                "minor", "structure",
                f"Missing expected section: '{section}'"
            )

    return review


def review_game_html(html_path: Path) -> OutputReview:
    """Review a generated HTML game for quality issues."""
    review = OutputReview()

    try:
        html = html_path.read_text(encoding="utf-8")
    except Exception as e:
        review.add_issue("critical", "file", f"Cannot read HTML: {e}")
        return review

    # Check: has proper HTML structure
    if "<script>" not in html and "<script " not in html:
        review.add_issue("critical", "structure", "No <script> tag — game won't run")

    if "<style>" not in html and "<style " not in html:
        review.add_issue("critical", "structure", "No <style> tag — game won't display")

    if "<body>" not in html and "<body " not in html:
        review.add_issue("major", "structure", "No <body> tag")

    # Check: has game elements
    game_indicators = [
        "addEventListener", "onclick", "function", "score",
        "querySelector", "getElementById",
    ]
    found = sum(1 for g in game_indicators if g in html)
    if found < 3:
        review.add_issue(
            "critical", "content",
            f"Only {found}/6 game indicators found — may not be interactive"
        )

    # Check: file size (too small = broken, too large = bloated)
    size_kb = len(html) / 1024
    if size_kb < 5:
        review.add_issue("critical", "content", f"Game is only {size_kb:.0f}KB — likely broken")
    elif size_kb > 200:
        review.add_issue("minor", "size", f"Game is {size_kb:.0f}KB — may load slowly")

    # Check: duplicate DOCTYPE
    if html.count("<!DOCTYPE") > 1:
        review.add_issue("major", "structure", "Duplicate <!DOCTYPE> tags")

    return review


def review_all_outputs(output_dir: Path) -> dict[str, OutputReview]:
    """Review all generated outputs in a directory."""
    reviews: dict[str, OutputReview] = {}

    for pptx in output_dir.glob("*.pptx"):
        reviews[pptx.name] = review_pptx(pptx)

    for docx in output_dir.glob("*handout*.docx"):
        reviews[docx.name] = review_docx(docx)

    for html in output_dir.glob("game_*.html"):
        reviews[html.name] = review_game_html(html)

    return reviews
