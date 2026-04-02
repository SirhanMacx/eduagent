"""File ingestion pipeline — extracts text, images, and URLs from teacher curriculum files.

Supported formats: PDF, DOCX, PPTX, TXT, MD, SMART Notebook (.notebook),
SMART Board (.xbk), ActivInspire (.flipchart), and ZIP archives.
"""

from __future__ import annotations

import csv
import io
import logging
import re
import shutil
import subprocess
import tempfile
import zipfile
from collections import Counter
from pathlib import Path
from typing import Optional

from clawed.asset_registry import ExtractedImage, ExtractedURL, ExtractionResult, extract_urls, extract_youtube_ids
from clawed.models import DocType, Document

logger = logging.getLogger(__name__)

# ── File-type detection ──────────────────────────────────────────────────

EXTENSION_MAP: dict[str, DocType] = {
    ".pdf": DocType.PDF,
    ".docx": DocType.DOCX,
    ".pptx": DocType.PPTX,
    ".txt": DocType.TXT,
    ".md": DocType.MD,
    ".notebook": DocType.NOTEBOOK,
    ".xbk": DocType.XBK,
    ".flipchart": DocType.FLIPCHART,
    # Added for v2.3.7 — common legacy/teacher formats
    ".doc": DocType.DOCX,
    ".ppt": DocType.PPTX,
    ".xls": DocType.TXT,
    ".xlsx": DocType.TXT,
    ".csv": DocType.TXT,
    ".rtf": DocType.TXT,
    ".html": DocType.TXT,
    ".htm": DocType.TXT,
    ".odt": DocType.TXT,
    ".odp": DocType.TXT,
}

SUPPORTED_EXTENSIONS = set(EXTENSION_MAP.keys())


def _detect_type(path: Path) -> DocType:
    return EXTENSION_MAP.get(path.suffix.lower(), DocType.UNKNOWN)


# ── URL extraction helpers ───────────────────────────────────────────────

_YT_RE = re.compile(
    r'(?:https?://)?(?:www\.|m\.)?(?:youtube\.com/watch\?v=|youtu\.be/|youtube\.com/embed/)'
    r'([a-zA-Z0-9_-]{11})'
)


def _extract_urls_from_text(text: str) -> list[ExtractedURL]:
    """Extract URLs from plain text, classifying YouTube vs other."""
    urls: list[ExtractedURL] = []
    seen: set[str] = set()
    for url in extract_urls(text):
        if url in seen:
            continue
        seen.add(url)
        yt_ids = extract_youtube_ids(url)
        if yt_ids:
            urls.append(ExtractedURL(
                url=f"https://youtube.com/watch?v={yt_ids[0]}",
                link_type="youtube",
            ))
        elif 'docs.google.com' in url or 'drive.google.com' in url:
            urls.append(ExtractedURL(url=url, link_type="google_doc"))
        else:
            urls.append(ExtractedURL(url=url, link_type="website"))
    return urls


# ── Individual extractors ────────────────────────────────────────────────


def _extract_pdf(path: Path) -> tuple[str, Optional[int]]:
    """Extract text from a PDF using PyMuPDF."""
    import fitz  # PyMuPDF

    doc = fitz.open(str(path))
    pages: list[str] = []
    for page in doc:
        text = page.get_text()
        if text.strip():
            pages.append(text)
    page_count = len(doc)
    doc.close()
    return "\n\n".join(pages), page_count


def _extract_pdf_rich(path: Path) -> ExtractionResult:
    """Extract text, images, and URLs from a PDF."""
    import fitz

    doc = fitz.open(str(path))
    pages: list[str] = []
    images: list[ExtractedImage] = []
    for page_num, page in enumerate(doc):
        text = page.get_text()
        slide_context = ""
        if text.strip():
            pages.append(text)
            slide_context = text.strip()
        for img_info in page.get_images():
            xref = img_info[0]
            try:
                base_image = doc.extract_image(xref)
                if base_image and base_image.get("image"):
                    img_bytes = base_image["image"]
                    ext = base_image.get("ext", "png")
                    if len(img_bytes) > 2000:  # Skip tiny icons
                        images.append(ExtractedImage(
                            image_bytes=img_bytes,
                            format=ext,
                            context_text=slide_context[:200] if slide_context else "",
                        ))
            except Exception as e:
                logger.debug("Failed to extract image xref %d: %s", xref, e)
    page_count = len(doc)
    doc.close()

    full_text = "\n\n".join(pages)
    urls = _extract_urls_from_text(full_text)

    return ExtractionResult(
        text=full_text,
        page_count=page_count,
        word_count=len(full_text.split()),
        urls=urls,
        images=images,
    )


def _extract_docx(path: Path) -> str:
    """Extract text from a DOCX file."""
    from docx import Document as DocxDocument

    doc = DocxDocument(str(path))
    paragraphs: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            paragraphs.append(para.text)
    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))
    return "\n\n".join(paragraphs)


def _extract_docx_rich(path: Path) -> ExtractionResult:
    """Extract text, images, and hyperlinks from a DOCX."""
    from docx import Document as DocxDocument

    doc = DocxDocument(str(path))
    paragraphs: list[str] = []
    urls: list[ExtractedURL] = []
    images: list[ExtractedImage] = []

    # Extract text
    for para in doc.paragraphs:
        if para.text.strip():
            paragraphs.append(para.text)

    # Extract table text
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))

    full_text = "\n\n".join(paragraphs)

    # Extract hyperlinks from document XML
    try:
        rels = doc.part.rels
        for rel in rels.values():
            if "hyperlink" in rel.reltype:
                url = rel._target  # noqa: SLF001
                if url and isinstance(url, str) and url.startswith("http"):
                    yt_ids = extract_youtube_ids(url)
                    if yt_ids:
                        urls.append(ExtractedURL(
                            url=f"https://youtube.com/watch?v={yt_ids[0]}",
                            link_type="youtube",
                        ))
                    else:
                        urls.append(ExtractedURL(url=url, link_type="website"))
    except (AttributeError, KeyError, ValueError) as e:
        logger.warning("Failed to extract DOCX hyperlinks from %s: %s", path.name, e)

    # Also find URLs in plain text
    for text_url in _extract_urls_from_text(full_text):
        if not any(u.url == text_url.url for u in urls):
            urls.append(text_url)

    # Extract inline images
    try:
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    img_part = rel.target_part
                    img_bytes = img_part.blob
                    ct = img_part.content_type or ""
                    fmt = "png" if "png" in ct else "jpeg"
                    if len(img_bytes) > 1000:  # skip tiny icons
                        images.append(ExtractedImage(
                            image_bytes=img_bytes,
                            format=fmt,
                            alt_text="",
                            context_text="",
                        ))
                except (AttributeError, OSError, ValueError) as e:
                    logger.debug("Failed to extract DOCX image: %s", e)
    except (AttributeError, KeyError) as e:
        logger.warning("Failed to enumerate DOCX image rels from %s: %s", path.name, e)

    return ExtractionResult(
        text=full_text,
        word_count=len(full_text.split()),
        urls=urls,
        images=images,
    )


def _extract_pptx(path: Path) -> str:
    """Extract text from a PPTX file."""
    from pptx import Presentation

    prs = Presentation(str(path))
    slides_text: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts: list[str] = [f"[Slide {i}]"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        parts.append(text)
        if len(parts) > 1:  # More than just the slide marker
            slides_text.append("\n".join(parts))
    return "\n\n".join(slides_text)


def _extract_pptx_rich(path: Path) -> ExtractionResult:
    """Extract text, images, and hyperlinks from a PPTX."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(path))
    slides_text: list[str] = []
    urls: list[ExtractedURL] = []
    images: list[ExtractedImage] = []
    seen_urls: set[str] = set()

    for i, slide in enumerate(prs.slides, 1):
        parts: list[str] = [f"[Slide {i}]"]
        slide_context = ""

        for shape in slide.shapes:
            # Extract text
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        parts.append(text)
                        slide_context = text[:100] if not slide_context else slide_context

                    # Extract hyperlinks from runs
                    for run in paragraph.runs:
                        try:
                            if run.hyperlink and run.hyperlink.address:
                                url = run.hyperlink.address
                                if url not in seen_urls:
                                    seen_urls.add(url)
                                    yt_ids = extract_youtube_ids(url)
                                    if yt_ids:
                                        urls.append(ExtractedURL(
                                            url=f"https://youtube.com/watch?v={yt_ids[0]}",
                                            link_type="youtube",
                                            context_text=run.text[:100],
                                        ))
                                    else:
                                        urls.append(ExtractedURL(
                                            url=url,
                                            link_type="website",
                                            context_text=run.text[:100],
                                        ))
                        except (AttributeError, ValueError) as e:
                            logger.debug("Failed to extract PPTX hyperlink: %s", e)

            # Extract images
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img = shape.image
                    ct = img.content_type or ""
                    # Skip EMF/WMF (not reusable in new PPTX easily)
                    if 'emf' in ct.lower() or 'wmf' in ct.lower():
                        continue
                    fmt = "png" if "png" in ct else "jpeg"
                    img_bytes = img.blob
                    if len(img_bytes) > 2000:  # skip tiny icons/bullets
                        images.append(ExtractedImage(
                            image_bytes=img_bytes,
                            format=fmt,
                            width=shape.width,
                            height=shape.height,
                            context_text=slide_context,
                            slide_number=i,
                        ))
            except (AttributeError, OSError, ValueError) as e:
                logger.debug("Failed to extract PPTX image on slide %d: %s", i, e)

        if len(parts) > 1:
            slides_text.append("\n".join(parts))

    full_text = "\n\n".join(slides_text)

    # Also find URLs in plain text
    for text_url in _extract_urls_from_text(full_text):
        if text_url.url not in seen_urls:
            urls.append(text_url)

    return ExtractionResult(
        text=full_text,
        slide_count=len(prs.slides),
        word_count=len(full_text.split()),
        urls=urls,
        images=images,
    )


def _extract_text(path: Path) -> str:
    """Extract text from a plain text or Markdown file."""
    return path.read_text(encoding="utf-8", errors="replace")


def _extract_notebook(path: Path) -> tuple[str, Optional[int]]:
    """Extract text from a SMART Notebook file via formats plugin."""
    from clawed.formats.notebook import extract_notebook

    return extract_notebook(path)


def _extract_xbk(path: Path) -> tuple[str, Optional[int]]:
    """Extract text from a SMART Board file via formats plugin."""
    from clawed.formats.xbk import extract_xbk

    return extract_xbk(path)


def _extract_flipchart(path: Path) -> tuple[str, Optional[int]]:
    """Extract text from an ActivInspire file via formats plugin."""
    from clawed.formats.flipchart import extract_flipchart

    return extract_flipchart(path)


# ── Legacy/additional format extractors (v2.3.7) ────────────────────────


def _extract_doc(path: Path) -> str:
    """Extract text from a legacy .doc file (binary Word format).

    Tries textutil (macOS), then catdoc/antiword, then raw text fallback.
    """
    # macOS built-in
    if shutil.which("textutil"):
        try:
            result = subprocess.run(
                ["textutil", "-convert", "txt", "-stdout", str(path)],
                capture_output=True, text=True, timeout=30,
            )
            if result.returncode == 0 and result.stdout.strip():
                return result.stdout.strip()
        except (OSError, subprocess.TimeoutExpired) as e:
            logger.warning("textutil failed on %s: %s", path.name, e)
    # catdoc fallback
    for tool in ("catdoc", "antiword"):
        if shutil.which(tool):
            try:
                result = subprocess.run(
                    [tool, str(path)],
                    capture_output=True, text=True, timeout=30,
                )
                if result.returncode == 0 and result.stdout.strip():
                    return result.stdout.strip()
            except (OSError, subprocess.TimeoutExpired) as e:
                logger.warning("%s failed on %s: %s", tool, path.name, e)
    # Raw text fallback (often garbled but better than nothing)
    try:
        raw = path.read_bytes()
        text = raw.decode("utf-8", errors="replace")
        # Strip null bytes and control chars
        text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", " ", text)
        text = re.sub(r"\s{3,}", "\n", text).strip()
        if len(text) > 50:
            return text
    except OSError as e:
        logger.warning("Failed to read .doc raw bytes from %s: %s", path.name, e)
    return ""


def _extract_ppt(path: Path) -> str:
    """Extract text from a legacy .ppt file (binary PowerPoint format).

    Tries textutil (macOS), then raw text extraction.
    """
    if shutil.which("textutil"):
        try:
            result = subprocess.run(
                ["textutil", "-convert", "txt", "-stdout", str(path)],
                capture_output=True, text=True, timeout=30,
            )
            if result.returncode == 0 and result.stdout.strip():
                return result.stdout.strip()
        except (OSError, subprocess.TimeoutExpired) as e:
            logger.warning("textutil failed on %s: %s", path.name, e)
    # Raw text fallback
    try:
        raw = path.read_bytes()
        text = raw.decode("utf-8", errors="replace")
        text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", " ", text)
        text = re.sub(r"\s{3,}", "\n", text).strip()
        if len(text) > 50:
            return text
    except OSError as e:
        logger.warning("Failed to read .ppt raw bytes from %s: %s", path.name, e)
    return ""


def _extract_csv(path: Path) -> str:
    """Extract text from a CSV file."""
    try:
        content = path.read_text(encoding="utf-8", errors="replace")
        reader = csv.reader(io.StringIO(content))
        rows = []
        for i, row in enumerate(reader):
            if i >= 5000:  # Cap at 5000 rows
                rows.append(f"... ({i} rows total, truncated)")
                break
            rows.append(" | ".join(row))
        return "\n".join(rows)
    except (OSError, csv.Error, ValueError) as e:
        logger.warning("Failed to extract CSV from %s: %s", path.name, e)
        return ""


def _extract_xlsx(path: Path) -> str:
    """Extract text from an .xlsx file using openpyxl."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True, data_only=True)
        rows: list[str] = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            rows.append(f"[Sheet: {sheet}]")
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= 5000:
                    rows.append("... (truncated at 5000 rows)")
                    break
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append(" | ".join(cells))
        wb.close()
        return "\n".join(rows)
    except ImportError:
        logger.debug("openpyxl not installed, cannot extract .xlsx: %s", path.name)
        return ""
    except Exception as exc:
        logger.debug("Failed to extract .xlsx %s: %s", path.name, exc)
        return ""


def _extract_xls(path: Path) -> str:
    """Extract text from a legacy .xls file.

    Tries xlrd (proper XLS parsing), then textutil (macOS), then falls back
    to attempting xlsx parsing (some .xls files are actually xlsx in disguise).
    """
    # Strategy 1: xlrd (proper XLS parser)
    try:
        import xlrd
        wb = xlrd.open_workbook(str(path))
        rows: list[str] = []
        for sheet in wb.sheet_names():
            ws = wb.sheet_by_name(sheet)
            rows.append(f"[Sheet: {sheet}]")
            for row_idx in range(min(ws.nrows, 5000)):
                cells = []
                for col_idx in range(ws.ncols):
                    cell = ws.cell(row_idx, col_idx)
                    val = str(cell.value) if cell.value is not None else ""
                    cells.append(val)
                if any(c.strip() for c in cells):
                    rows.append(" | ".join(cells))
            if ws.nrows > 5000:
                rows.append("... (truncated at 5000 rows)")
        wb.release_resources()
        text = "\n".join(rows)
        if text.strip():
            return text
    except ImportError:
        logger.debug("xlrd not installed; cannot parse .xls natively: %s", path.name)
    except Exception as exc:
        logger.debug("xlrd failed on %s: %s", path.name, exc)

    # Strategy 2: textutil on macOS
    if shutil.which("textutil"):
        try:
            result = subprocess.run(
                ["textutil", "-convert", "txt", "-stdout", str(path)],
                capture_output=True, text=True, timeout=30,
            )
            if result.returncode == 0 and result.stdout.strip():
                return result.stdout.strip()
        except (OSError, subprocess.TimeoutExpired) as e:
            logger.warning("textutil failed on %s: %s", path.name, e)

    # Strategy 3: Try openpyxl (some .xls files are actually xlsx format)
    try:
        text = _extract_xlsx(path)
        if text.strip():
            logger.debug("Parsed .xls as xlsx format: %s", path.name)
            return text
    except Exception:
        pass

    logger.debug(
        "Could not extract .xls file %s — install xlrd for full .xls support",
        path.name,
    )
    return ""


def _extract_rtf(path: Path) -> str:
    """Extract text from an RTF file.

    Tries striprtf library first, then textutil (macOS), then regex fallback.
    """
    try:
        from striprtf.striprtf import rtf_to_text
        raw = path.read_text(encoding="utf-8", errors="replace")
        return rtf_to_text(raw).strip()
    except ImportError:
        pass
    except (OSError, ValueError) as e:
        logger.warning("striprtf failed on %s: %s", path.name, e)
    # macOS textutil
    if shutil.which("textutil"):
        try:
            result = subprocess.run(
                ["textutil", "-convert", "txt", "-stdout", str(path)],
                capture_output=True, text=True, timeout=30,
            )
            if result.returncode == 0 and result.stdout.strip():
                return result.stdout.strip()
        except (OSError, subprocess.TimeoutExpired) as e:
            logger.warning("textutil failed on RTF %s: %s", path.name, e)
    # Regex fallback: strip RTF control codes
    try:
        raw = path.read_text(encoding="utf-8", errors="replace")
        # Remove RTF control words and braces
        text = re.sub(r"\\[a-z]+\d*\s?", " ", raw)
        text = re.sub(r"[{}]", "", text)
        text = re.sub(r"\s{2,}", " ", text).strip()
        if len(text) > 50:
            return text
    except OSError as e:
        logger.warning("Failed to read RTF file %s: %s", path.name, e)
    return ""


def _extract_html_file(path: Path) -> str:
    """Extract text from an HTML file by stripping tags."""
    try:
        from html.parser import HTMLParser

        class _TagStripper(HTMLParser):
            def __init__(self):
                super().__init__()
                self.parts: list[str] = []
                self._skip = False

            def handle_starttag(self, tag, attrs):
                if tag in ("script", "style"):
                    self._skip = True

            def handle_endtag(self, tag):
                if tag in ("script", "style"):
                    self._skip = False

            def handle_data(self, data):
                if not self._skip:
                    self.parts.append(data)

        raw = path.read_text(encoding="utf-8", errors="replace")
        stripper = _TagStripper()
        stripper.feed(raw)
        text = " ".join(stripper.parts)
        text = re.sub(r"\s{2,}", " ", text).strip()
        return text
    except (OSError, ValueError) as e:
        logger.warning("Failed to extract HTML from %s: %s", path.name, e)
        return ""


def _odf_iter_text(element) -> str:
    """Recursively collect all text from an XML element and its children."""
    parts: list[str] = []
    if element.text:
        parts.append(element.text)
    for child in element:
        parts.append(_odf_iter_text(child))
        if child.tail:
            parts.append(child.tail)
    return "".join(parts)


def _parse_odf_content_xml(content_bytes: bytes) -> list[str]:
    """Parse ODF content.xml using ElementTree to preserve structure.

    Handles ODF namespaces for text:h (headings), text:p (paragraphs),
    text:list-item (list items), and draw:page (presentation slides).
    Returns a list of text blocks with structure preserved.
    """
    import xml.etree.ElementTree as ET

    ns = {
        "text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0",
        "draw": "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0",
        "presentation": "urn:oasis:names:tc:opendocument:xmlns:presentation:1.0",
        "office": "urn:oasis:names:tc:opendocument:xmlns:office:1.0",
    }

    blocks: list[str] = []

    try:
        root = ET.fromstring(content_bytes)
    except ET.ParseError as exc:
        logger.debug("Failed to parse ODF content.xml: %s", exc)
        return blocks

    # Collect headings (text:h) with a prefix marker
    for heading in root.iter(f"{{{ns['text']}}}h"):
        text = _odf_iter_text(heading).strip()
        if text:
            blocks.append(f"## {text}")

    # Collect paragraphs (text:p) — these live inside text:section,
    # text:list-item, draw:text-box, etc.  We walk the whole tree.
    for para in root.iter(f"{{{ns['text']}}}p"):
        text = _odf_iter_text(para).strip()
        if text:
            blocks.append(text)

    # Collect list items (text:list-item) — deduplicate with paragraphs
    # by checking if the text is already captured.  List items contain
    # text:p children, so we only add the bullet marker.
    seen = set(blocks)
    for li in root.iter(f"{{{ns['text']}}}list-item"):
        text = _odf_iter_text(li).strip()
        if text and text not in seen:
            blocks.append(f"  - {text}")
            seen.add(text)

    return blocks


def _extract_odt(path: Path) -> str:
    """Extract text from an OpenDocument Text (.odt) file.

    ODT is a ZIP containing content.xml.  Parses XML properly to preserve
    paragraph structure, headings, and lists.
    """
    try:
        with zipfile.ZipFile(str(path), "r") as zf:
            if "content.xml" not in zf.namelist():
                return ""
            content_bytes = zf.read("content.xml")
            blocks = _parse_odf_content_xml(content_bytes)
            if not blocks:
                # Fallback: decode and strip tags (better than nothing)
                raw = content_bytes.decode("utf-8", errors="replace")
                text = re.sub(r"<[^>]+>", " ", raw)
                text = re.sub(r"\s{2,}", " ", text).strip()
                return text
            return "\n\n".join(blocks)
    except Exception as exc:
        logger.debug("Failed to extract .odt %s: %s", path.name, exc)
        return ""


def _extract_odp(path: Path) -> str:
    """Extract text from an OpenDocument Presentation (.odp) file.

    ODP is a ZIP containing content.xml with <draw:page> elements for slides.
    Parses XML properly to preserve slide structure.
    """
    import xml.etree.ElementTree as ET

    ns_draw = "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0"
    ns_text = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"

    try:
        with zipfile.ZipFile(str(path), "r") as zf:
            if "content.xml" not in zf.namelist():
                return ""
            content_bytes = zf.read("content.xml")

            try:
                root = ET.fromstring(content_bytes)
            except ET.ParseError:
                # Fallback to naive stripping
                raw = content_bytes.decode("utf-8", errors="replace")
                text = re.sub(r"<[^>]+>", " ", raw)
                text = re.sub(r"\s{2,}", " ", text).strip()
                return text

            slides_text: list[str] = []
            slide_num = 0
            for page in root.iter(f"{{{ns_draw}}}page"):
                slide_num += 1
                page_name = page.get(f"{{{ns_draw}}}name", f"Slide {slide_num}")
                parts: list[str] = [f"[Slide {slide_num}: {page_name}]"]

                # Collect all text:p and text:h elements within this slide
                for elem in page.iter():
                    tag = elem.tag
                    if tag == f"{{{ns_text}}}h" or tag == f"{{{ns_text}}}p":
                        text = _odf_iter_text(elem).strip()
                        if text:
                            parts.append(text)

                if len(parts) > 1:
                    slides_text.append("\n".join(parts))

            if not slides_text:
                # No draw:page elements found — try generic text extraction
                blocks = _parse_odf_content_xml(content_bytes)
                return "\n\n".join(blocks) if blocks else ""

            return "\n\n".join(slides_text)
    except Exception as exc:
        logger.debug("Failed to extract .odp %s: %s", path.name, exc)
        return ""


# Extension-based extractors for new formats that need special handling
# (keyed by suffix, returns (text, page_count_or_None))
_EXTENSION_EXTRACTORS: dict[str, object] = {
    ".doc": lambda p: (_extract_doc(p), None),
    ".ppt": lambda p: (_extract_ppt(p), None),
    ".csv": lambda p: (_extract_csv(p), None),
    ".xlsx": lambda p: (_extract_xlsx(p), None),
    ".xls": lambda p: (_extract_xls(p), None),
    ".rtf": lambda p: (_extract_rtf(p), None),
    ".html": lambda p: (_extract_html_file(p), None),
    ".htm": lambda p: (_extract_html_file(p), None),
    ".odt": lambda p: (_extract_odt(p), None),
    ".odp": lambda p: (_extract_odp(p), None),
}


# ── Topic tag extraction ──────────────────────────────────────────────


# Words too common to be useful as topic tags
_TAG_STOPWORDS = frozenset({
    "the", "a", "an", "and", "or", "but", "in", "on", "at", "to", "for",
    "of", "with", "by", "from", "is", "are", "was", "were", "be", "been",
    "have", "has", "had", "do", "does", "did", "will", "would", "could",
    "should", "may", "might", "shall", "can", "this", "that", "these",
    "those", "it", "its", "not", "no", "so", "if", "as", "up", "out",
    "about", "into", "through", "during", "before", "after", "above",
    "below", "between", "under", "again", "further", "then", "once",
    "here", "there", "when", "where", "why", "how", "all", "each",
    "every", "both", "few", "more", "most", "other", "some", "such",
    "than", "too", "very", "just", "also", "new", "old", "first",
    "last", "next", "per", "page", "slide", "sheet", "file", "name",
    "class", "grade", "date", "unit", "week", "day", "period", "mr",
    "mrs", "ms", "dr", "student", "students", "teacher",
})


def _extract_topic_tags(path: Path, content: str) -> list[str]:
    """Extract topic tags from filename, metadata, and content.

    Sources (in priority order):
    1. Filename — split on underscores/hyphens, filter stopwords
    2. Document title/heading — first heading-like line in content
    3. Content analysis — capitalized phrases and repeated terms in first 500 chars

    Returns a deduplicated list of lowercase tags (max 10).
    """
    tags: list[str] = []
    seen: set[str] = set()

    def _add(tag: str) -> None:
        t = tag.strip().lower()
        # Filter out stopwords, too-short, or purely numeric tags
        if t and t not in seen and t not in _TAG_STOPWORDS and len(t) > 2 and not t.isdigit():
            seen.add(t)
            tags.append(t)

    # ── Source 1: Filename ──────────────────────────────────────────
    stem = path.stem
    # Split on underscores, hyphens, dots, and camelCase boundaries
    filename_parts = re.split(r"[_\-.\s]+", stem)
    for part in filename_parts:
        # Split camelCase
        sub_parts = re.sub(r"([a-z])([A-Z])", r"\1 \2", part).split()
        for sp in sub_parts:
            _add(sp)

    # ── Source 2: First heading or title line ───────────────────────
    # Look for lines that look like titles/headings (short, capitalized)
    for line in content.split("\n")[:20]:
        line = line.strip()
        if not line:
            continue
        # Markdown heading
        if line.startswith("#"):
            heading_text = line.lstrip("#").strip()
            for word in heading_text.split():
                _add(word)
            break
        # All-caps or title-case short line (likely a heading)
        if len(line) < 80 and (line.isupper() or line.istitle()):
            for word in line.split():
                _add(word)
            break

    # ── Source 3: Content analysis (first 500 chars) ───────────────
    snippet = content[:500]

    # Capitalized multi-word phrases (potential proper nouns / topic names)
    # Use [ \t]+ instead of \s+ to avoid matching across newlines
    cap_phrases = re.findall(r"\b([A-Z][a-z]+(?:[ \t]+[A-Z][a-z]+)+)\b", snippet)
    for phrase in cap_phrases[:5]:
        _add(phrase)

    # Repeated terms (words appearing 3+ times suggest topicality)
    word_counts: Counter[str] = Counter()
    for word in re.findall(r"\b[a-zA-Z]{3,}\b", snippet):
        word_counts[word.lower()] += 1
    for word, count in word_counts.most_common(10):
        if count >= 2 and word not in _TAG_STOPWORDS:
            _add(word)

    return tags[:10]


# ── Corpus contribution helper ─────────────────────────────────────────


def _contribute_to_corpus(doc: "Document") -> None:
    """Contribute an extracted document to the corpus if it looks like teaching material.

    Called during ingestion so that high-quality curriculum materials are
    available as few-shot examples for future lesson generation.
    Only contributes documents that contain teaching-related keywords.
    Failures are logged but never block ingestion.
    """
    try:
        from clawed.corpus import contribute_example
    except ImportError:
        return

    content_lower = doc.content.lower()

    # Detect lesson-plan-like documents
    is_lesson = any(kw in content_lower for kw in [
        "objective", "swbat", "do now", "aim", "warm up",
        "direct instruction", "guided practice", "exit ticket",
        "homework", "materials needed", "lesson plan",
    ])
    is_unit = any(kw in content_lower for kw in [
        "unit plan", "essential question", "enduring understanding",
        "unit overview", "unit goals", "pacing guide",
    ])

    if not (is_lesson or is_unit):
        return

    content_type = "lesson_plan" if is_lesson else "unit_plan"

    # Infer subject from content heuristics (best-effort)
    subject = "general"
    subject_signals = {
        "social studies": ["social studies", "history", "civilization", "government"],
        "science": ["science", "biology", "chemistry", "physics", "hypothesis", "experiment"],
        "math": ["math", "algebra", "geometry", "equation", "calculate"],
        "ela": ["reading", "writing", "literature", "grammar", "essay", "author"],
    }
    for subj, keywords in subject_signals.items():
        if any(kw in content_lower for kw in keywords):
            subject = subj
            break

    try:
        contribute_example(
            content_type=content_type,
            subject=subject,
            grade_level="9-12",  # default; refined by persona later
            content={"title": doc.title, "text": doc.content[:3000], "source_file": doc.title},
            topic=doc.title,
            quality_score=3.5,  # moderate default — teacher feedback raises it
            source="ingest",
        )
        logger.debug("Contributed '%s' to corpus as %s (%s)", doc.title, content_type, subject)
    except Exception as exc:
        logger.debug("Corpus contribution failed for '%s': %s", doc.title, exc)


# ── Dispatch ─────────────────────────────────────────────────────────────

EXTRACTORS = {
    DocType.PDF: lambda p: _extract_pdf(p),
    DocType.DOCX: lambda p: (_extract_docx(p), None),
    DocType.PPTX: lambda p: (_extract_pptx(p), None),
    DocType.TXT: lambda p: (_extract_text(p), None),
    DocType.MD: lambda p: (_extract_text(p), None),
    DocType.NOTEBOOK: lambda p: _extract_notebook(p),
    DocType.XBK: lambda p: _extract_xbk(p),
    DocType.FLIPCHART: lambda p: _extract_flipchart(p),
}

# Rich extractors return ExtractionResult with images, URLs, metadata
RICH_EXTRACTORS = {
    DocType.PPTX: _extract_pptx_rich,
    DocType.DOCX: _extract_docx_rich,
    DocType.PDF: _extract_pdf_rich,
}


def _extract_single(path: Path) -> Optional[Document]:
    """Extract a single file into a Document, or None if unsupported/empty."""
    doc_type = _detect_type(path)
    if doc_type == DocType.UNKNOWN:
        logger.warning("Skipping unsupported format: %s", path.name)
        return None

    # Check extension-based extractors first (for legacy formats like .doc,
    # .ppt, .csv, etc. that map to a generic DocType but need special handling)
    ext_extractor = _EXTENSION_EXTRACTORS.get(path.suffix.lower())
    extractor = ext_extractor or EXTRACTORS.get(doc_type)
    if not extractor:
        logger.warning("No extractor for %s format: %s", doc_type.value, path.name)
        return None

    try:
        result = extractor(path)
    except Exception as exc:
        logger.warning("Failed to parse %s: %s", path.name, exc)
        return None
    if isinstance(result, tuple):
        content, page_count = result
    else:
        content, page_count = result, None

    if not content or not content.strip():
        return None

    cleaned_content = content.strip()
    topic_tags = _extract_topic_tags(path, cleaned_content)

    doc = Document(
        title=path.stem.replace("_", " ").replace("-", " ").title(),
        content=cleaned_content,
        doc_type=doc_type,
        source_path=str(path),
        page_count=page_count,
        tags=topic_tags,
    )

    # Contribute to corpus (best-effort, never blocks ingestion)
    try:
        _contribute_to_corpus(doc)
    except Exception as exc:
        logger.debug("Corpus contribution skipped for %s: %s", path.name, exc)

    return doc


def extract_rich(path: Path) -> Optional[ExtractionResult]:
    """Extract rich metadata (images, URLs, slide counts) from a file.

    Returns None if the file type doesn't support rich extraction or if
    extraction fails. The caller should fall back to _extract_single().
    """
    doc_type = _detect_type(path)
    extractor = RICH_EXTRACTORS.get(doc_type)
    if not extractor:
        return None
    try:
        return extractor(path)
    except Exception as exc:
        logger.debug("Rich extraction failed for %s: %s", path.name, exc)
        return None


# ── Format summary ──────────────────────────────────────────────────────


def _collect_files(path: Path) -> list[Path]:
    """Collect all supported files from a directory, sorted by modification
    time (newest first).

    Skips Office lock files (~$*) and macOS resource forks (._*).
    Sorting by mtime means the most recent materials come first, which is
    important when a cap is applied.
    """
    files = [
        f for f in path.rglob("*")
        if f.is_file()
        and f.suffix.lower() in SUPPORTED_EXTENSIONS
        and not f.name.startswith(("~$", "._"))
    ]
    # Sort newest first so caps keep the most recent materials
    files.sort(key=lambda f: f.stat().st_mtime, reverse=True)
    return files


def _format_summary(files: list[Path]) -> str:
    """Build a human-readable summary of found file types.

    Example: "Found 47 PDFs, 23 DOCX, 12 PPTX, 3 TXT — analyzing..."
    """
    counts: Counter[str] = Counter()
    for f in files:
        ext = f.suffix.lower().lstrip(".")
        counts[ext.upper()] += 1

    parts = [f"{count} {fmt}" for fmt, count in sorted(counts.items(), key=lambda x: -x[1])]
    return f"Found {', '.join(parts)} — analyzing..."


# ── Public API ───────────────────────────────────────────────────────────


def scan_directory(path: Path) -> tuple[list[Path], str]:
    """Scan a directory and return supported files with a format summary.

    Returns (file_list, summary_string).
    """
    if not path.is_dir():
        raise FileNotFoundError(f"Directory not found: {path}")

    files = _collect_files(path)
    summary = _format_summary(files) if files else "No supported files found."
    return files, summary


def ingest_directory(
    path: Path,
    *,
    max_files: int = 0,
    progress_callback=None,
) -> list[Document]:
    """Recursively scan a directory and extract all supported documents.

    Args:
        path: Directory to scan.
        max_files: Maximum number of files to process (0 = unlimited).
            Files are already sorted newest-first by ``_collect_files``.
        progress_callback: Optional callable(current, total) for progress updates.
    """
    documents: list[Document] = []
    if not path.is_dir():
        raise FileNotFoundError(f"Directory not found: {path}")

    files = _collect_files(path)
    if max_files > 0:
        files = files[:max_files]
    total = len(files)

    for i, file_path in enumerate(files):
        doc = _extract_single(file_path)
        if doc:
            documents.append(doc)
        if progress_callback:
            progress_callback(i + 1, total)

    return documents


MAX_UNZIP_SIZE = 500 * 1024 * 1024  # 500 MB


def ingest_zip(path: Path, *, progress_callback=None) -> list[Document]:
    """Extract a ZIP file to a temp directory, then ingest its contents."""
    if not path.is_file() or path.suffix.lower() != ".zip":
        raise ValueError(f"Not a ZIP file: {path}")

    with tempfile.TemporaryDirectory() as tmp_dir:
        with zipfile.ZipFile(str(path), "r") as zf:
            total_size = sum(info.file_size for info in zf.infolist())
            if total_size > MAX_UNZIP_SIZE:
                logger.warning(
                    "ZIP file too large when decompressed (%d bytes), skipping",
                    total_size,
                )
                return []
            zf.extractall(tmp_dir)
        return ingest_directory(Path(tmp_dir), progress_callback=progress_callback)


def ingest_path(
    path: Path,
    *,
    dry_run: bool = False,
    max_files: int = 0,
    progress_callback=None,
) -> list[Document]:
    """Smart ingestion: accepts a directory, ZIP file, or single file.

    Args:
        path: File or directory to ingest.
        dry_run: If True, scan and log what would be processed but don't extract.
        max_files: Maximum number of files to process (0 = unlimited).
        progress_callback: Optional callable(current, total) for progress updates.
    """
    path = Path(path).expanduser().resolve()

    if path.is_dir():
        if dry_run:
            files, summary = scan_directory(path)
            logger.info(summary)
            return _dry_run_results(files)
        return ingest_directory(
            path, max_files=max_files, progress_callback=progress_callback,
        )
    elif path.is_file() and path.suffix.lower() == ".zip":
        if dry_run:
            # For ZIP dry-run, list contents without extracting
            with tempfile.TemporaryDirectory() as tmp_dir:
                with zipfile.ZipFile(str(path), "r") as zf:
                    total_size = sum(info.file_size for info in zf.infolist())
                    if total_size > MAX_UNZIP_SIZE:
                        logger.warning(
                            "ZIP file too large when decompressed (%d bytes), skipping",
                            total_size,
                        )
                        return []
                    zf.extractall(tmp_dir)
                files, summary = scan_directory(Path(tmp_dir))
                logger.info(summary)
                return _dry_run_results(files)
        return ingest_zip(path, progress_callback=progress_callback)
    elif path.is_file():
        if dry_run:
            return _dry_run_results([path])
        doc = _extract_single(path)
        return [doc] if doc else []
    else:
        raise FileNotFoundError(f"Path not found: {path}")


def _dry_run_results(files: list[Path]) -> list[Document]:
    """Create placeholder Documents for dry-run mode (no content extraction)."""
    results: list[Document] = []
    for f in files:
        doc_type = _detect_type(f)
        if doc_type != DocType.UNKNOWN:
            results.append(Document(
                title=f.stem.replace("_", " ").replace("-", " ").title(),
                content="[dry-run: content not extracted]",
                doc_type=doc_type,
                source_path=str(f),
            ))
    return results
