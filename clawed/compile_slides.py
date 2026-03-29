"""PPTX slide compiler for MasterContent.

Compiles a MasterContent object into a classroom-ready PowerPoint slide deck.
Slide order: title → vocabulary → instruction sections → source analysis →
station overview → exit ticket.
No LLM calls — pure mechanical compilation.
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import TYPE_CHECKING

if TYPE_CHECKING:
    from clawed.master_content import MasterContent

logger = logging.getLogger(__name__)


def _hex_to_rgb(hex_color: str):
    """Convert a 6-char hex string to pptx RGBColor."""
    from pptx.dml.color import RGBColor

    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _add_slide(prs, layout_idx: int = 6):
    """Add a blank slide and return it."""
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def _textbox(slide, left, top, width, height, text: str,
             font_size: int = 18, bold: bool = False,
             hex_color: str = "222222", align_center: bool = False,
             italic: bool = False, word_wrap: bool = True):
    """Add a textbox to a slide and return the shape."""
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    if align_center:
        p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _hex_to_rgb(hex_color)
    run.font.name = "Calibri"
    return tb


def _bullet_textbox(slide, left, top, width, height,
                    items: list[str], font_size: int = 16,
                    hex_color: str = "333333"):
    """Add a textbox with one paragraph per bullet item."""
    from pptx.util import Pt

    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"\u2022  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = _hex_to_rgb(hex_color)
        run.font.name = "Calibri"

    return tb


def _embed_image(slide, image_spec: str, images: dict[str, Path],
                 left, top, width, height) -> bool:
    """Embed a pre-fetched image on the slide. Returns True if embedded."""
    if not image_spec:
        return False
    path = images.get(image_spec)
    if path and Path(path).exists():
        try:
            slide.shapes.add_picture(str(path), left, top, width, height)
            return True
        except Exception as exc:
            logger.debug("Could not embed slide image %r: %s", image_spec, exc)
    return False


async def compile_slides(
    master: "MasterContent",
    images: dict[str, Path],
    output_dir: Path,
) -> Path:
    """Compile a classroom-ready PPTX from a MasterContent object.

    Slide order:
        1. Title slide (title, subject, grade, objective)
        2. Vocabulary slide(s) (up to 5 terms per slide)
        3. One slide per InstructionSection (heading, key points, content)
        4. Source analysis slides (one per primary source)
        5. Station overview (if stations exist)
        6. Exit ticket slide (questions only, no answers)

    Args:
        master: The MasterContent source-of-truth object.
        images: Mapping of image_spec strings to local file Paths.
        output_dir: Directory where the .pptx file will be written.

    Returns:
        Path to the generated .pptx file.
    """
    from pptx import Presentation
    from pptx.util import Emu, Inches, Pt

    from clawed.io import safe_filename

    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    W = prs.slide_width  # noqa: N806
    _slide_h = prs.slide_height  # noqa: F841, N806

    # Palette (neutral academic)
    C_TITLE_BG = "1F3864"    # dark navy  # noqa: N806
    C_SECTION_BG = "2E75B6"  # medium blue  # noqa: N806
    C_WHITE = "FFFFFF"  # noqa: N806
    C_DARK = "222222"  # noqa: N806
    C_ACCENT = "BDD7EE"      # light blue  # noqa: N806

    def _set_bg(slide, hex_color: str) -> None:
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb(hex_color)

    # ── 1. Title slide ────────────────────────────────────────────────

    slide = _add_slide(prs)
    _set_bg(slide, C_TITLE_BG)

    # Main title
    _textbox(
        slide,
        left=Inches(1.0), top=Inches(1.5),
        width=W - Inches(2.0), height=Inches(1.8),
        text=master.title,
        font_size=40, bold=True,
        hex_color=C_WHITE, align_center=True,
    )
    # Subject / grade / duration
    meta = f"{master.subject}  |  Grade {master.grade_level}  |  {master.duration_minutes} min"
    _textbox(
        slide,
        left=Inches(1.0), top=Inches(3.4),
        width=W - Inches(2.0), height=Inches(0.5),
        text=meta,
        font_size=18, bold=False,
        hex_color=C_ACCENT, align_center=True,
    )
    # Objective
    _textbox(
        slide,
        left=Inches(1.0), top=Inches(4.1),
        width=W - Inches(2.0), height=Inches(1.4),
        text=f"Objective: {master.objective}",
        font_size=16,
        hex_color=C_WHITE, align_center=True,
    )

    # ── 2. Vocabulary slide(s) ────────────────────────────────────────

    TERMS_PER_SLIDE = 5  # noqa: N806
    if master.vocabulary:
        vocab_chunks = [
            master.vocabulary[i: i + TERMS_PER_SLIDE]
            for i in range(0, len(master.vocabulary), TERMS_PER_SLIDE)
        ]
        for chunk_idx, chunk in enumerate(vocab_chunks):
            slide = _add_slide(prs)
            heading_label = "Vocabulary" if len(vocab_chunks) == 1 else f"Vocabulary ({chunk_idx + 1})"
            # Section header bar
            bar = slide.shapes.add_textbox(
                Inches(0), Inches(0), W, Inches(0.8),
            )
            bar_tf = bar.text_frame
            bar_tf.word_wrap = False
            bar_p = bar_tf.paragraphs[0]
            from pptx.enum.text import PP_ALIGN
            bar_p.alignment = PP_ALIGN.LEFT
            bar_run = bar_p.add_run()
            bar_run.text = f"  {heading_label}"
            bar_run.font.size = Pt(24)
            bar_run.font.bold = True
            bar_run.font.color.rgb = _hex_to_rgb(C_WHITE)
            bar_run.font.name = "Calibri"
            # Fill bar background via shape fill (add a rectangle behind)
            from pptx.util import Emu
            rect = slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
                Emu(0), Emu(0), W, Inches(0.8),
            )
            rect.fill.solid()
            rect.fill.fore_color.rgb = _hex_to_rgb(C_SECTION_BG)
            rect.line.fill.background()
            # Move bar text on top by re-adding it after the rect
            slide.shapes._spTree.remove(bar._element)
            slide.shapes._spTree.append(bar._element)

            # Term table: two columns (term | definition + context)
            top_offset = Inches(1.0)
            row_height = Inches(1.0)
            for entry in chunk:
                # Term box
                _textbox(
                    slide,
                    left=Inches(0.3), top=top_offset,
                    width=Inches(3.0), height=row_height,
                    text=entry.term,
                    font_size=16, bold=True,
                    hex_color=C_DARK,
                )
                # Definition + context box
                defn_text = entry.definition
                if entry.context_sentence:
                    defn_text += f"\n\u201c{entry.context_sentence}\u201d"
                _textbox(
                    slide,
                    left=Inches(3.5), top=top_offset,
                    width=Inches(9.0), height=row_height,
                    text=defn_text,
                    font_size=14,
                    hex_color=C_DARK,
                )
                top_offset += row_height + Inches(0.05)

    # ── 3. One slide per InstructionSection ───────────────────────────

    for section in master.direct_instruction:
        slide = _add_slide(prs)

        # Header bar (reuse pattern)
        rect = slide.shapes.add_shape(1, Emu(0), Emu(0), W, Inches(0.8))
        rect.fill.solid()
        rect.fill.fore_color.rgb = _hex_to_rgb(C_SECTION_BG)
        rect.line.fill.background()
        _textbox(
            slide,
            left=Inches(0.2), top=Inches(0.05),
            width=W - Inches(0.4), height=Inches(0.7),
            text=section.heading,
            font_size=22, bold=True,
            hex_color=C_WHITE,
        )

        has_image = bool(section.image_spec and section.image_spec in images)
        content_width = W - Inches(5.0) if has_image else W - Inches(0.6)
        content_left = Inches(0.3)

        # Key points (bullets)
        if section.key_points:
            _bullet_textbox(
                slide,
                left=content_left, top=Inches(1.0),
                width=content_width, height=Inches(2.5),
                items=section.key_points,
                font_size=18,
                hex_color=C_DARK,
            )

        # Content summary (truncated at 400 chars to keep slides readable)
        summary = section.content[:400] + ("…" if len(section.content) > 400 else "")
        _textbox(
            slide,
            left=content_left, top=Inches(3.6),
            width=content_width, height=Inches(2.8),
            text=summary,
            font_size=14,
            hex_color="444444",
        )

        # Image (right side)
        if has_image:
            _embed_image(
                slide, section.image_spec, images,
                left=W - Inches(4.6), top=Inches(1.0),
                width=Inches(4.3), height=Inches(5.0),
            )

    # ── 4. Source analysis slides ─────────────────────────────────────

    for ps in master.primary_sources:
        slide = _add_slide(prs)

        # Header
        rect = slide.shapes.add_shape(1, Emu(0), Emu(0), W, Inches(0.8))
        rect.fill.solid()
        rect.fill.fore_color.rgb = _hex_to_rgb("C55A11")  # burnt orange for sources
        rect.line.fill.background()
        _textbox(
            slide,
            left=Inches(0.2), top=Inches(0.05),
            width=W - Inches(0.4), height=Inches(0.7),
            text=f"Source Analysis: {ps.title}",
            font_size=20, bold=True,
            hex_color=C_WHITE,
        )

        has_image = bool(ps.image_spec and ps.image_spec in images)
        text_width = W - Inches(5.0) if has_image else W - Inches(0.6)

        # Attribution + type
        _textbox(
            slide,
            left=Inches(0.3), top=Inches(0.9),
            width=text_width, height=Inches(0.4),
            text=f"{ps.source_type.replace('_', ' ').title()}  |  {ps.attribution}",
            font_size=12, italic=True,
            hex_color="666666",
        )

        # Excerpt (first 300 chars)
        excerpt = ps.content_text[:300] + ("…" if len(ps.content_text) > 300 else "")
        _textbox(
            slide,
            left=Inches(0.3), top=Inches(1.4),
            width=text_width, height=Inches(2.5),
            text=f'"{excerpt}"',
            font_size=14, italic=True,
            hex_color=C_DARK,
        )

        # Scaffolding questions
        if ps.scaffolding_questions:
            _bullet_textbox(
                slide,
                left=Inches(0.3), top=Inches(4.0),
                width=text_width, height=Inches(2.5),
                items=ps.scaffolding_questions,
                font_size=15,
                hex_color=C_DARK,
            )

        # Image
        if has_image:
            _embed_image(
                slide, ps.image_spec, images,
                left=W - Inches(4.6), top=Inches(1.0),
                width=Inches(4.3), height=Inches(5.0),
            )

    # ── 5. Station overview (if stations exist) ───────────────────────

    if master.stations:
        slide = _add_slide(prs)

        rect = slide.shapes.add_shape(1, Emu(0), Emu(0), W, Inches(0.8))
        rect.fill.solid()
        rect.fill.fore_color.rgb = _hex_to_rgb("375623")  # dark green
        rect.line.fill.background()
        _textbox(
            slide,
            left=Inches(0.2), top=Inches(0.05),
            width=W - Inches(0.4), height=Inches(0.7),
            text="Learning Stations",
            font_size=24, bold=True,
            hex_color=C_WHITE,
        )

        # List each station title + task
        top_offset = Inches(1.0)
        col_w = (W - Inches(0.6)) / max(len(master.stations), 1)
        for i, station in enumerate(master.stations):
            _textbox(
                slide,
                left=Inches(0.3) + col_w * i,
                top=top_offset,
                width=col_w - Inches(0.1),
                height=Inches(1.0),
                text=station.title,
                font_size=16, bold=True,
                hex_color=C_DARK,
            )
            _textbox(
                slide,
                left=Inches(0.3) + col_w * i,
                top=top_offset + Inches(1.0),
                width=col_w - Inches(0.1),
                height=Inches(4.5),
                text=station.student_directions,
                font_size=13,
                hex_color="444444",
            )

    # ── 6. Exit ticket slide (questions only, no answers) ─────────────

    if master.exit_ticket:
        slide = _add_slide(prs)

        rect = slide.shapes.add_shape(1, Emu(0), Emu(0), W, Inches(0.8))
        rect.fill.solid()
        rect.fill.fore_color.rgb = _hex_to_rgb("7030A0")  # purple
        rect.line.fill.background()
        _textbox(
            slide,
            left=Inches(0.2), top=Inches(0.05),
            width=W - Inches(0.4), height=Inches(0.7),
            text="Exit Ticket",
            font_size=24, bold=True,
            hex_color=C_WHITE,
        )

        top_offset = Inches(1.0)
        for i, sq in enumerate(master.exit_ticket, 1):
            has_image = bool(sq.stimulus_image_spec and sq.stimulus_image_spec in images)
            q_width = W - Inches(5.0) if has_image else W - Inches(0.6)
            # Stimulus
            _textbox(
                slide,
                left=Inches(0.3), top=top_offset,
                width=q_width, height=Inches(0.8),
                text=f"Stimulus: {sq.stimulus[:200]}",
                font_size=13, italic=True,
                hex_color="444444",
            )
            # Question
            _textbox(
                slide,
                left=Inches(0.3), top=top_offset + Inches(0.85),
                width=q_width, height=Inches(0.8),
                text=f"Q{i}: {sq.question}",
                font_size=16, bold=True,
                hex_color=C_DARK,
            )
            if has_image:
                _embed_image(
                    slide, sq.stimulus_image_spec, images,
                    left=W - Inches(4.6), top=top_offset,
                    width=Inches(4.3), height=Inches(2.0),
                )
            top_offset += Inches(2.0)

    # ── Save ──────────────────────────────────────────────────────────

    safe = safe_filename(master.title)
    out_path = output_dir / f"{safe}_slides.pptx"
    prs.save(str(out_path))
    logger.info("Slides saved to %s", out_path)
    return out_path
