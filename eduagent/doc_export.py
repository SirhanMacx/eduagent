"""Document export — generate PPTX, DOCX, and PDF from lesson plans.

Uses python-pptx for PowerPoint and python-docx for Word documents,
both of which are already project dependencies.  PDF is generated
from the DOCX via reportlab (also a dependency).
"""

from __future__ import annotations

import asyncio
import logging
from datetime import date
from pathlib import Path
from typing import TYPE_CHECKING, Any, Optional

if TYPE_CHECKING:
    from eduagent.models import DailyLesson, TeacherPersona

logger = logging.getLogger(__name__)


# ── Enhanced color theme definitions ──────────────────────────────────

_COLOR_THEMES: dict[str, dict[str, str]] = {
    "history": {
        "primary": "8B4513",       # Saddle brown
        "secondary": "DAA520",     # Goldenrod
        "accent": "F5E6CC",        # Cream
        "bg_dark": "2C1810",       # Dark brown
        "bg_light": "FFF8F0",      # Warm cream
        "text_dark": "1A1A1A",
        "text_light": "FFFFFF",
    },
    "social studies": {
        "primary": "8B4513",
        "secondary": "DAA520",
        "accent": "F5E6CC",
        "bg_dark": "2C1810",
        "bg_light": "FFF8F0",
        "text_dark": "1A1A1A",
        "text_light": "FFFFFF",
    },
    "science": {
        "primary": "1B5E20",       # Dark green
        "secondary": "43A047",     # Medium green
        "accent": "E8F5E9",        # Mint
        "bg_dark": "0D3311",       # Deep green
        "bg_light": "F0F8F5",      # Light mint
        "text_dark": "1A1A1A",
        "text_light": "FFFFFF",
    },
    "biology": {
        "primary": "1B5E20",
        "secondary": "43A047",
        "accent": "E8F5E9",
        "bg_dark": "0D3311",
        "bg_light": "F0F8F5",
        "text_dark": "1A1A1A",
        "text_light": "FFFFFF",
    },
    "chemistry": {
        "primary": "1B5E20",
        "secondary": "43A047",
        "accent": "E8F5E9",
        "bg_dark": "0D3311",
        "bg_light": "F0F8F5",
        "text_dark": "1A1A1A",
        "text_light": "FFFFFF",
    },
    "physics": {
        "primary": "1B5E20",
        "secondary": "43A047",
        "accent": "E8F5E9",
        "bg_dark": "0D3311",
        "bg_light": "F0F8F5",
        "text_dark": "1A1A1A",
        "text_light": "FFFFFF",
    },
    "math": {
        "primary": "1565C0",       # Blue
        "secondary": "42A5F5",     # Light blue
        "accent": "E3F2FD",        # Ice blue
        "bg_dark": "0D2137",       # Navy
        "bg_light": "F0F4FA",      # Soft blue
        "text_dark": "1A1A1A",
        "text_light": "FFFFFF",
    },
    "mathematics": {
        "primary": "1565C0",
        "secondary": "42A5F5",
        "accent": "E3F2FD",
        "bg_dark": "0D2137",
        "bg_light": "F0F4FA",
        "text_dark": "1A1A1A",
        "text_light": "FFFFFF",
    },
    "algebra": {
        "primary": "1565C0",
        "secondary": "42A5F5",
        "accent": "E3F2FD",
        "bg_dark": "0D2137",
        "bg_light": "F0F4FA",
        "text_dark": "1A1A1A",
        "text_light": "FFFFFF",
    },
    "ela": {
        "primary": "6A1B9A",       # Purple
        "secondary": "AB47BC",     # Light purple
        "accent": "F3E5F5",        # Lavender
        "bg_dark": "2A0845",       # Deep purple
        "bg_light": "F5F0FA",      # Light lavender
        "text_dark": "1A1A1A",
        "text_light": "FFFFFF",
    },
    "english": {
        "primary": "6A1B9A",
        "secondary": "AB47BC",
        "accent": "F3E5F5",
        "bg_dark": "2A0845",
        "bg_light": "F5F0FA",
        "text_dark": "1A1A1A",
        "text_light": "FFFFFF",
    },
    "language arts": {
        "primary": "6A1B9A",
        "secondary": "AB47BC",
        "accent": "F3E5F5",
        "bg_dark": "2A0845",
        "bg_light": "F5F0FA",
        "text_dark": "1A1A1A",
        "text_light": "FFFFFF",
    },
}

_DEFAULT_THEME: dict[str, str] = {
    "primary": "1A365D",       # Professional navy
    "secondary": "3182CE",     # Blue
    "accent": "EBF8FF",        # Light blue
    "bg_dark": "0A1628",       # Dark navy
    "bg_light": "F0F4FA",      # Soft blue-white
    "text_dark": "1A1A1A",
    "text_light": "FFFFFF",
}


def get_color_theme(subject: str) -> dict[str, str]:
    """Return the color theme dict for a subject, or the default."""
    return _COLOR_THEMES.get(subject.strip().lower(), _DEFAULT_THEME)


# ── PowerPoint export ──────────────────────────────────────────────────


def _detect_subject(persona: "TeacherPersona") -> str:
    """Best-effort subject detection from persona fields."""
    subj = (persona.subject_area or "").strip().lower()
    if subj:
        return subj
    return ""


def _hex_to_rgb(hex_str: str):
    """Convert a hex color string to an RGBColor."""
    from pptx.dml.color import RGBColor

    return RGBColor(int(hex_str[:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def _add_shape_fill(shape, hex_color: str) -> None:
    """Fill a shape with a solid color."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(hex_color)


def _set_text_props(run, font_size_pt: int, hex_color: str, bold: bool = False):
    """Set font properties on a text run."""
    from pptx.util import Pt

    run.font.size = Pt(font_size_pt)
    run.font.color.rgb = _hex_to_rgb(hex_color)
    run.font.bold = bold
    run.font.name = "Calibri"


def _try_fetch_images(topics: list[tuple[str, str]], subject: str) -> dict[str, Optional[Path]]:
    """Attempt to fetch images for multiple topics. Non-blocking, short timeout.

    Returns a dict mapping key -> Path | None.
    """
    from eduagent.slide_images import fetch_slide_image

    results: dict[str, Optional[Path]] = {}

    async def _fetch_all():
        for topic, key in topics:
            try:
                path = await asyncio.wait_for(
                    fetch_slide_image(topic, subject=subject),
                    timeout=5.0,
                )
                results[key] = path
            except Exception:
                results[key] = None

    try:
        try:
            loop = asyncio.get_running_loop()
        except RuntimeError:
            loop = None

        if loop and loop.is_running():
            import concurrent.futures
            with concurrent.futures.ThreadPoolExecutor(max_workers=1) as pool:
                future = pool.submit(asyncio.run, _fetch_all())
                future.result(timeout=30)
        else:
            asyncio.run(_fetch_all())
    except Exception as e:
        logger.debug("Image fetching failed: %s", e)

    return results


def _try_fetch_content_images(
    items: list[tuple[str, str, str]],
    subject: str,
    max_images: int = 4,
) -> dict[str, Optional[Path]]:
    """Fetch images based on slide *content* text, not just the lesson title.

    Each item is ``(content_text, fallback_topic, key)``.  Stops after
    ``max_images`` successful fetches to avoid slowing down generation.

    Returns a dict mapping key -> Path | None.
    """
    from eduagent.slide_images import fetch_content_image, fetch_slide_image

    results: dict[str, Optional[Path]] = {}

    async def _fetch_all():
        found = 0
        for content_text, fallback_topic, key in items:
            if found >= max_images:
                results[key] = None
                continue
            try:
                if content_text:
                    path = await asyncio.wait_for(
                        fetch_content_image(
                            content_text,
                            subject=subject,
                            fallback_topic=fallback_topic,
                        ),
                        timeout=5.0,
                    )
                else:
                    # Title slide: use topic-based search
                    path = await asyncio.wait_for(
                        fetch_slide_image(fallback_topic, subject=subject),
                        timeout=5.0,
                    )
                results[key] = path
                if path:
                    found += 1
            except Exception:
                results[key] = None

    try:
        try:
            loop = asyncio.get_running_loop()
        except RuntimeError:
            loop = None

        if loop and loop.is_running():
            import concurrent.futures
            with concurrent.futures.ThreadPoolExecutor(max_workers=1) as pool:
                future = pool.submit(asyncio.run, _fetch_all())
                future.result(timeout=30)
        else:
            asyncio.run(_fetch_all())
    except Exception as e:
        logger.debug("Content image fetching failed: %s", e)

    return results


def export_lesson_pptx(
    lesson: "DailyLesson",
    persona: "TeacherPersona",
    output_dir: Path | None = None,
) -> Path:
    """Generate a professional PowerPoint presentation from a lesson plan.

    Produces a polished, school-board-ready slide deck with:
    - Subject-themed color palette
    - Large readable fonts for projection on classroom screens
    - Academic images (LOC, Wikimedia, Unsplash) where available
    - Clean, modern, minimal design -- NOT busy or cluttered
    - Consistent visual language throughout

    Returns the path to the saved .pptx file.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Inches, Pt

    subject = _detect_subject(persona)
    theme = get_color_theme(subject)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # ── Try to fetch images asynchronously (per-slide content queries) ─
    # Each slide gets its OWN image query based on that slide's content.
    # Max 3-4 images per deck to avoid slow generation.
    # Images on: Title (bg), Do Now (accent), Direct Instruction (sidebar)
    # No images on: Objectives, Guided Practice, Exit Ticket, Closing
    image_items: list[tuple[str, str, str]] = [
        # (content_text, fallback_topic, key)
        ("", lesson.title, "title"),  # title slide: topic-based
    ]
    if lesson.do_now:
        image_items.append((lesson.do_now, lesson.title, "do_now"))
    if lesson.direct_instruction:
        image_items.append(
            (lesson.direct_instruction, lesson.title, "instruction"),
        )

    images = _try_fetch_content_images(image_items, subject, max_images=3)

    # ── Shared layout helpers ─────────────────────────────────────────
    slide_num = [0]

    def _next_slide():
        slide_num[0] += 1
        layout = prs.slide_layouts[6]  # blank layout
        return prs.slides.add_slide(layout)

    def _add_footer(slide, num: int):
        left = slide_w - Inches(1.5)
        top = slide_h - Inches(0.45)
        tb = slide.shapes.add_textbox(left, top, Inches(1.2), Inches(0.3))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = str(num)
        _set_text_props(run, 10, "999999")

    def _bar(slide, left, top, width, height, hex_color: str):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height,
        )
        shape.line.fill.background()
        _add_shape_fill(shape, hex_color)
        return shape

    def _rounded_card(slide, left, top, width, height, hex_color: str):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
        )
        shape.line.fill.background()
        _add_shape_fill(shape, hex_color)
        # Subtle shadow
        try:
            shape.shadow.inherit = False
            shadow = shape.shadow
            shadow.visible = True
            shadow.blur_radius = Pt(4)
            shadow.distance = Pt(2)
        except Exception:
            pass
        return shape

    def _white_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb("FFFFFF")

    def _tinted_bg(slide, hex_color: str):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb(hex_color)

    def _add_bg_image(slide, image_path: Path, overlay_alpha: str = "30000"):
        """Full-bleed background image with dark gradient overlay."""
        pic = slide.shapes.add_picture(
            str(image_path), Emu(0), Emu(0), slide_w, slide_h,
        )
        sp = pic._element
        sp.getparent().remove(sp)
        slide.shapes._spTree.insert(2, sp)

        from pptx.oxml.ns import qn

        overlay = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), slide_w, slide_h,
        )
        overlay.line.fill.background()
        fill_obj = overlay.fill
        fill_obj.solid()
        fill_obj.fore_color.rgb = _hex_to_rgb(theme.get("bg_dark", theme["primary"]))
        solid_fill = overlay._element.spPr.solidFill
        if solid_fill is not None:
            srgb = solid_fill.find(qn("a:srgbClr"))
            if srgb is None:
                srgb = solid_fill.find(qn("a:sysClr"))
            if srgb is not None:
                from lxml import etree
                alpha = etree.SubElement(srgb, qn("a:alpha"))
                alpha.set("val", overlay_alpha)

    def _add_sidebar_image(slide, image_path: Path, caption: str = ""):
        img_left = int(slide_w * 0.65)
        img_width = int(slide_w * 0.33)
        img_top = Inches(1.4)
        img_height = int(slide_h - Inches(2.4))
        try:
            slide.shapes.add_picture(
                str(image_path), img_left, img_top, img_width, img_height,
            )
            if caption:
                cap_tb = slide.shapes.add_textbox(
                    img_left, img_top + img_height + Inches(0.05),
                    img_width, Inches(0.35),
                )
                cap_p = cap_tb.text_frame.paragraphs[0]
                cap_p.alignment = PP_ALIGN.CENTER
                cap_run = cap_p.add_run()
                cap_run.text = caption
                _set_text_props(cap_run, 9, "888888")
                cap_run.font.italic = True
        except Exception:
            pass

    def _add_accent_image(slide, image_path: Path, caption: str = ""):
        """Add a smaller accent image in the bottom-right area."""
        img_left = slide_w - Inches(4.0)
        img_top = slide_h - Inches(3.2)
        img_width = Inches(3.5)
        img_height = Inches(2.2)
        try:
            slide.shapes.add_picture(
                str(image_path), img_left, img_top, img_width, img_height,
            )
            if caption:
                cap_tb = slide.shapes.add_textbox(
                    img_left, img_top + img_height + Inches(0.05),
                    img_width, Inches(0.3),
                )
                cap_p = cap_tb.text_frame.paragraphs[0]
                cap_p.alignment = PP_ALIGN.CENTER
                cap_run = cap_p.add_run()
                cap_run.text = caption
                _set_text_props(cap_run, 9, "888888")
                cap_run.font.italic = True
        except Exception:
            pass

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 1: TITLE
    # Full-bleed image (or solid primary bg) with lesson title overlay
    # ═══════════════════════════════════════════════════════════════════
    slide = _next_slide()

    title_image = images.get("title")
    if title_image:
        _add_bg_image(slide, title_image, overlay_alpha="30000")  # 70% dark overlay
    else:
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb(theme.get("bg_dark", theme["primary"]))

    # Lesson title -- 44pt white bold, centered
    tb = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.0), slide_w - Inches(3.0), Inches(2.5),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = lesson.title
    _set_text_props(run, 44, theme["text_light"], bold=True)

    # Subtitle: "Lesson N | Teacher Name | Date" -- 20pt white
    tb2 = slide.shapes.add_textbox(
        Inches(1.5), Inches(4.8), slide_w - Inches(3.0), Inches(1.0),
    )
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = (
        f"Lesson {lesson.lesson_number}  |  "
        f"{persona.name or 'Teacher'}  |  "
        f"{date.today().strftime('%B %d, %Y')}"
    )
    _set_text_props(run2, 20, "DDDDDD")

    # Bottom accent bar in subject color
    _bar(slide, Emu(0), slide_h - Inches(0.15), slide_w, Inches(0.15), theme["secondary"])

    _add_footer(slide, slide_num[0])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 2: OBJECTIVES
    # Clean white bg, left accent bar, SWBAT in colored card
    # ═══════════════════════════════════════════════════════════════════
    slide = _next_slide()
    _white_bg(slide)

    # Left accent bar (4px, subject color)
    _bar(slide, Inches(0.5), Inches(0.8), Inches(0.06), Inches(5.5), theme["primary"])

    # "Today's Objectives" header -- 36pt, subject color
    tb = slide.shapes.add_textbox(
        Inches(1.0), Inches(0.8), slide_w - Inches(2.0), Inches(1.0),
    )
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Today's Objectives"
    _set_text_props(run, 36, theme["primary"], bold=True)

    # SWBAT card -- colored rounded rectangle
    _rounded_card(
        slide,
        Inches(1.0), Inches(2.2),
        slide_w - Inches(2.0), Inches(1.8),
        theme["accent"],
    )
    tb = slide.shapes.add_textbox(
        Inches(1.4), Inches(2.5), slide_w - Inches(2.8), Inches(1.2),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.space_after = Pt(12)
    run = p.add_run()
    run.text = f"SWBAT: {lesson.objective}"
    _set_text_props(run, 24, theme["text_dark"], bold=True)

    # Standards below in 16pt gray
    if lesson.standards:
        tb = slide.shapes.add_textbox(
            Inches(1.0), Inches(4.4), slide_w - Inches(2.0), Inches(2.5),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Standards Addressed"
        _set_text_props(run, 16, "888888", bold=True)

        for std in lesson.standards[:5]:
            p = tf.add_paragraph()
            p.space_before = Pt(6)
            run = p.add_run()
            run.text = f"  {std}"
            _set_text_props(run, 16, "666666")

    _add_footer(slide, slide_num[0])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 3: DO NOW / WARM-UP
    # Light tint bg, "Do Now" badge, question text, timer
    # ═══════════════════════════════════════════════════════════════════
    if lesson.do_now:
        slide = _next_slide()
        _tinted_bg(slide, theme["accent"])

        # "Do Now" badge -- rounded rectangle, subject color, white text
        badge = _rounded_card(
            slide,
            Inches(0.8), Inches(0.6),
            Inches(2.5), Inches(0.7),
            theme["primary"],
        )
        badge_tf = badge.text_frame
        badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = badge_tf.paragraphs[0].add_run()
        run.text = "Do Now"
        _set_text_props(run, 22, theme["text_light"], bold=True)

        # Question / prompt text -- 28pt, dark, good line spacing
        # Do Now uses accent image (small, bottom-right) not sidebar
        do_now_img = images.get("do_now")

        tb = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.8), slide_w - Inches(2.0), Inches(3.2),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.line_spacing = Pt(38)
        run = p.add_run()
        run.text = lesson.do_now
        _set_text_props(run, 28, theme["text_dark"])

        if do_now_img:
            from eduagent.slide_images import _extract_key_concepts
            concepts = _extract_key_concepts(lesson.do_now)
            caption = ", ".join(concepts[:2]) if concepts else ""
            _add_accent_image(slide, do_now_img, caption=caption)

        # Timer indicator
        minutes = lesson.time_estimates.get("do_now", 5)
        tb_timer = slide.shapes.add_textbox(
            slide_w - Inches(2.5), slide_h - Inches(1.0),
            Inches(2.0), Inches(0.5),
        )
        p_t = tb_timer.text_frame.paragraphs[0]
        p_t.alignment = PP_ALIGN.RIGHT
        run_t = p_t.add_run()
        run_t.text = f"{minutes} minutes"
        _set_text_props(run_t, 16, "888888")

        _add_footer(slide, slide_num[0])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDES 4+: DIRECT INSTRUCTION
    # Numbered section badge, key concept heading, side image
    # ═══════════════════════════════════════════════════════════════════
    if lesson.direct_instruction:
        text = lesson.direct_instruction
        chunks = _split_text(text, max_len=550) if len(text) > 600 else [text]

        for i, chunk in enumerate(chunks, 1):
            slide = _next_slide()
            _white_bg(slide)

            # Numbered circle badge (subject color)
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.8), Inches(0.6),
                Inches(0.8), Inches(0.8),
            )
            circle.line.fill.background()
            _add_shape_fill(circle, theme["primary"])
            circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = circle.text_frame.paragraphs[0].add_run()
            run.text = str(i)
            _set_text_props(run, 24, theme["text_light"], bold=True)

            # Heading -- 32pt bold
            suffix = f"  ({i} of {len(chunks)})" if len(chunks) > 1 else ""
            tb = slide.shapes.add_textbox(
                Inches(2.0), Inches(0.65),
                slide_w - Inches(3.0), Inches(0.8),
            )
            p = tb.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = f"Direct Instruction{suffix}"
            _set_text_props(run, 32, theme["text_dark"], bold=True)

            # Side accent bar
            _bar(slide, Inches(0.6), Inches(1.7), Inches(0.06), Inches(5.0), theme["secondary"])

            # Sidebar image (right 35%) on first chunk, with caption
            img_path = images.get("instruction") if i == 1 else None
            text_width = slide_w - Inches(2.0)
            if img_path:
                text_width = int(slide_w * 0.60)
                from eduagent.slide_images import _extract_key_concepts
                concepts = _extract_key_concepts(chunk)
                caption = ", ".join(concepts[:2]) if concepts else ""
                _add_sidebar_image(slide, img_path, caption=caption)

            # Body text -- 20pt with 1.5x line spacing
            tb = slide.shapes.add_textbox(
                Inches(1.0), Inches(1.7), text_width, Inches(5.0),
            )
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.line_spacing = Pt(30)
            run = p.add_run()
            run.text = chunk
            _set_text_props(run, 20, theme["text_dark"])

            _add_footer(slide, slide_num[0])

    # ═══════════════════════════════════════════════════════════════════
    # GUIDED PRACTICE -- "Your Turn" header
    # ═══════════════════════════════════════════════════════════════════
    if lesson.guided_practice:
        slide = _next_slide()
        _tinted_bg(slide, "F5F5F5")  # very light gray

        # "Your Turn" badge
        badge = _rounded_card(
            slide,
            Inches(0.8), Inches(0.6),
            Inches(3.0), Inches(0.7),
            theme["secondary"],
        )
        badge_tf = badge.text_frame
        badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = badge_tf.paragraphs[0].add_run()
        run.text = "Your Turn"
        _set_text_props(run, 22, theme["text_light"], bold=True)

        # No image on Guided Practice -- keep clean for readability

        # Activity instructions -- 24pt
        tb = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.8), slide_w - Inches(2.0), Inches(4.5),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.line_spacing = Pt(34)
        run = p.add_run()
        run.text = lesson.guided_practice
        _set_text_props(run, 24, theme["text_dark"])

        # Time estimate
        minutes = lesson.time_estimates.get("guided_practice", 15)
        tb_time = slide.shapes.add_textbox(
            slide_w - Inches(2.5), slide_h - Inches(1.0),
            Inches(2.0), Inches(0.5),
        )
        p_t = tb_time.text_frame.paragraphs[0]
        p_t.alignment = PP_ALIGN.RIGHT
        run_t = p_t.add_run()
        run_t.text = f"{minutes} minutes"
        _set_text_props(run_t, 16, "888888")

        _add_footer(slide, slide_num[0])

    # ═══════════════════════════════════════════════════════════════════
    # INDEPENDENT WORK -- "Independent Practice" header
    # ═══════════════════════════════════════════════════════════════════
    if lesson.independent_work:
        slide = _next_slide()
        _white_bg(slide)

        # "Independent Practice" badge
        badge = _rounded_card(
            slide,
            Inches(0.8), Inches(0.6),
            Inches(4.5), Inches(0.7),
            theme["primary"],
        )
        badge_tf = badge.text_frame
        badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = badge_tf.paragraphs[0].add_run()
        run.text = "Independent Practice"
        _set_text_props(run, 22, theme["text_light"], bold=True)

        # Left accent bar
        _bar(slide, Inches(0.6), Inches(1.7), Inches(0.06), Inches(5.0), theme["secondary"])

        # No image on Independent Work -- keep clean for readability

        # Task description -- 24pt
        tb = slide.shapes.add_textbox(
            Inches(1.0), Inches(1.8), slide_w - Inches(2.0), Inches(4.5),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.line_spacing = Pt(34)
        run = p.add_run()
        run.text = lesson.independent_work
        _set_text_props(run, 24, theme["text_dark"])

        # Time estimate in corner
        minutes = lesson.time_estimates.get("independent_work", 10)
        _rounded_card(
            slide,
            slide_w - Inches(3.0), slide_h - Inches(1.2),
            Inches(2.2), Inches(0.6),
            theme["accent"],
        )
        tb_time = slide.shapes.add_textbox(
            slide_w - Inches(3.0), slide_h - Inches(1.2),
            Inches(2.2), Inches(0.6),
        )
        p_t = tb_time.text_frame.paragraphs[0]
        p_t.alignment = PP_ALIGN.CENTER
        run_t = p_t.add_run()
        run_t.text = f"{minutes} minutes"
        _set_text_props(run_t, 16, theme["text_dark"])

        _add_footer(slide, slide_num[0])

    # ═══════════════════════════════════════════════════════════════════
    # EXIT TICKET -- distinctive assessment design
    # Colored banner header, numbered question cards with shadow
    # ═══════════════════════════════════════════════════════════════════
    if lesson.exit_ticket:
        slide = _next_slide()
        _white_bg(slide)

        # Full-width colored banner for "Exit Ticket"
        _bar(slide, Emu(0), Emu(0), slide_w, Inches(1.2), theme["primary"])
        tb = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.15), slide_w - Inches(1.6), Inches(0.9),
        )
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "Exit Ticket"
        _set_text_props(run, 36, theme["text_light"], bold=True)

        # Numbered question cards
        q_top = Inches(1.6)
        for i, q in enumerate(lesson.exit_ticket, 1):
            # Card background (rounded rectangle with shadow)
            _rounded_card(
                slide,
                Inches(1.0), q_top,
                slide_w - Inches(2.0), Inches(1.1),
                theme["accent"],
            )

            # Number circle
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(1.3), q_top + Inches(0.15),
                Inches(0.7), Inches(0.7),
            )
            circle.line.fill.background()
            _add_shape_fill(circle, theme["secondary"])
            circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = circle.text_frame.paragraphs[0].add_run()
            run.text = str(i)
            _set_text_props(run, 22, theme["text_light"], bold=True)

            # Question text
            tb = slide.shapes.add_textbox(
                Inches(2.3), q_top + Inches(0.15),
                slide_w - Inches(3.8), Inches(0.8),
            )
            tf = tb.text_frame
            tf.word_wrap = True
            run = tf.paragraphs[0].add_run()
            run.text = q.question
            _set_text_props(run, 20, theme["text_dark"])

            q_top += Inches(1.3)

        # "Turn in before you leave" footer note
        tb = slide.shapes.add_textbox(
            Inches(0.8), slide_h - Inches(0.9),
            slide_w - Inches(1.6), Inches(0.4),
        )
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "Turn in before you leave"
        _set_text_props(run, 14, "888888")

        _add_footer(slide, slide_num[0])

    # ═══════════════════════════════════════════════════════════════════
    # CLOSING SLIDE -- Homework or Key Takeaway
    # ═══════════════════════════════════════════════════════════════════
    slide = _next_slide()
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(theme.get("bg_dark", theme["primary"]))

    if lesson.homework:
        # "Tonight's Homework" with details
        tb = slide.shapes.add_textbox(
            Inches(1.5), Inches(1.5), slide_w - Inches(3.0), Inches(1.5),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "Tonight's Homework"
        _set_text_props(run, 36, theme["text_light"], bold=True)

        _bar(
            slide, Inches(1.5), Inches(3.2),
            Inches(4.0), Inches(0.06), theme["secondary"],
        )

        tb = slide.shapes.add_textbox(
            Inches(1.5), Inches(3.6), slide_w - Inches(3.0), Inches(3.0),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.line_spacing = Pt(30)
        run = p.add_run()
        run.text = lesson.homework
        _set_text_props(run, 22, "DDDDDD")
    else:
        # "Key Takeaway" or "Questions?"
        tb = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.5), slide_w - Inches(3.0), Inches(2.0),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "Questions?"
        _set_text_props(run, 48, theme["text_light"], bold=True)

    # Teacher name and subject in footer
    teacher_subject = subject.title() if subject else "Education"
    tb = slide.shapes.add_textbox(
        Inches(1.0), slide_h - Inches(1.0),
        slide_w - Inches(2.0), Inches(0.5),
    )
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"{persona.name or 'Teacher'}  |  {teacher_subject}"
    _set_text_props(run, 14, "888888")

    # Watermark
    tb_wm = slide.shapes.add_textbox(
        Inches(1.0), slide_h - Inches(0.55),
        slide_w - Inches(2.0), Inches(0.3),
    )
    p_wm = tb_wm.text_frame.paragraphs[0]
    p_wm.alignment = PP_ALIGN.CENTER
    run_wm = p_wm.add_run()
    run_wm.text = "Generated by EDUagent"
    _set_text_props(run_wm, 10, "666666")

    _add_footer(slide, slide_num[0])

    # ── Save ──────────────────────────────────────────────────────────
    out = _resolve_output(output_dir, lesson, ".pptx")
    prs.save(str(out))
    return out


# ── Word document export ───────────────────────────────────────────────


def export_lesson_docx(
    lesson: "DailyLesson",
    persona: "TeacherPersona",
    output_dir: Path | None = None,
) -> Path:
    """Generate a Word document from a lesson plan with embedded academic images.

    Returns the path to the saved .docx file.
    """
    from docx import Document
    from docx.shared import Pt

    doc = Document()

    # Resolve subject for image searches
    subject = persona.subject_area or ""

    # Title
    doc.add_heading(lesson.title, level=0)
    doc.add_paragraph(
        f"Teacher: {persona.name or 'Teacher'}  |  "
        f"Lesson {lesson.lesson_number}  |  "
        f"{date.today().strftime('%B %d, %Y')}"
    )

    # Try to add a header image relevant to the lesson content
    _docx_add_content_image(
        doc,
        content_text=lesson.title + " " + lesson.objective,
        fallback_topic=lesson.title,
        subject=subject,
        width_inches=5.5,
    )

    # Standards table
    if lesson.standards:
        doc.add_heading("Standards Addressed", level=2)
        table = doc.add_table(rows=1, cols=1)
        table.style = "Light Grid Accent 1"
        hdr = table.rows[0].cells
        hdr[0].text = "Standard"
        for std in lesson.standards:
            row = table.add_row().cells
            row[0].text = std

    # Objective
    doc.add_heading("Objective (SWBAT)", level=2)
    doc.add_paragraph(lesson.objective)

    # Materials
    if lesson.materials_needed:
        doc.add_heading("Materials Needed", level=2)
        for m in lesson.materials_needed:
            doc.add_paragraph(m, style="List Bullet")

    # Lesson Sections — add relevant images to instruction sections
    sections = [
        ("Do Now / Warm-Up", lesson.do_now, False),
        ("Direct Instruction", lesson.direct_instruction, True),
        ("Guided Practice", lesson.guided_practice, False),
        ("Independent Work", lesson.independent_work, False),
    ]
    for heading, content, add_img in sections:
        if content:
            time_key = heading.lower().replace(" / warm-up", "").replace(" ", "_")
            minutes = lesson.time_estimates.get(time_key, "")
            time_label = f" ({minutes} min)" if minutes else ""
            doc.add_heading(f"{heading}{time_label}", level=2)
            doc.add_paragraph(content)
            # Add a content-specific image to the direct instruction section
            if add_img:
                _docx_add_content_image(
                    doc,
                    content_text=content,
                    fallback_topic=lesson.title,
                    subject=subject,
                    width_inches=4.0,
                )

    # Exit Ticket
    if lesson.exit_ticket:
        doc.add_heading("Exit Ticket", level=2)
        for i, q in enumerate(lesson.exit_ticket, 1):
            doc.add_paragraph(f"{i}. {q.question}")

    # Differentiation
    diff = lesson.differentiation
    if diff:
        doc.add_heading("Differentiation", level=2)
        if diff.struggling:
            doc.add_paragraph(f"Struggling learners: {diff.struggling}")
        if diff.advanced:
            doc.add_paragraph(f"Advanced learners: {diff.advanced}")
        if diff.ell:
            doc.add_paragraph(f"ELL support: {diff.ell}")

    # Homework
    if lesson.homework:
        doc.add_heading("Homework", level=2)
        doc.add_paragraph(lesson.homework)

    # Footer
    doc.add_paragraph("")
    footer = doc.add_paragraph("Generated by EDUagent")
    footer.runs[0].font.size = Pt(8)

    out = _resolve_output(output_dir, lesson, ".docx")
    doc.save(str(out))
    return out


def _docx_add_image(
    doc: "Any",  # docx.Document
    topic: str,
    subject: str,
    width_inches: float = 5.0,
    caption: str = "",
) -> bool:
    """Try to fetch and embed an academic image into the DOCX.

    Fails silently -- handout works fine without images.
    Returns True if an image was successfully added.
    """
    try:
        import asyncio

        from docx.shared import Inches, Pt

        from eduagent.slide_images import fetch_slide_image

        img_path = asyncio.run(fetch_slide_image(topic, subject=subject))
        if img_path and img_path.exists():
            doc.add_picture(str(img_path), width=Inches(width_inches))
            # Center the image
            last_para = doc.paragraphs[-1]
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            last_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            # Add caption if provided
            if caption:
                cap_para = doc.add_paragraph()
                cap_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                cap_run = cap_para.add_run(caption)
                cap_run.font.size = Pt(9)
                cap_run.font.italic = True
                cap_run.font.name = "Calibri"
                from docx.shared import RGBColor
                cap_run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
            return True
    except Exception:
        pass  # Images are a bonus, not a requirement
    return False


def _docx_add_content_image(
    doc: "Any",  # docx.Document
    content_text: str,
    fallback_topic: str,
    subject: str,
    width_inches: float = 3.0,
) -> bool:
    """Fetch and embed an image based on content text, with auto-caption.

    Uses ``_extract_key_concepts`` to find the best search query, then
    adds a captioned image.  Returns True if an image was added.
    """
    try:
        import asyncio

        from docx.shared import Inches, Pt

        from eduagent.slide_images import _extract_key_concepts, fetch_content_image

        img_path = asyncio.run(
            fetch_content_image(
                content_text,
                subject=subject,
                fallback_topic=fallback_topic,
            )
        )
        if img_path and img_path.exists():
            doc.add_picture(str(img_path), width=Inches(width_inches))
            last_para = doc.paragraphs[-1]
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            last_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            # Build caption from key concepts
            concepts = _extract_key_concepts(content_text)
            caption = ", ".join(concepts[:2]) if concepts else ""
            if caption:
                cap_para = doc.add_paragraph()
                cap_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                cap_run = cap_para.add_run(caption)
                cap_run.font.size = Pt(9)
                cap_run.font.italic = True
                cap_run.font.name = "Calibri"
                from docx.shared import RGBColor
                cap_run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
            return True
    except Exception:
        pass
    return False


# ── PDF export ─────────────────────────────────────────────────────────


def export_lesson_pdf(
    lesson: "DailyLesson",
    persona: "TeacherPersona",
    output_dir: Path | None = None,
) -> Path:
    """Generate a PDF from a lesson plan via reportlab.

    Returns the path to the saved .pdf file.
    """
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    out = _resolve_output(output_dir, lesson, ".pdf")

    doc = SimpleDocTemplate(
        str(out),
        pagesize=letter,
        topMargin=0.75 * inch,
        bottomMargin=0.75 * inch,
        leftMargin=1 * inch,
        rightMargin=1 * inch,
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "LessonTitle",
        parent=styles["Title"],
        fontSize=20,
        spaceAfter=6,
    )
    heading_style = ParagraphStyle(
        "SectionHeading",
        parent=styles["Heading2"],
        fontSize=14,
        spaceBefore=12,
        spaceAfter=4,
    )
    body_style = ParagraphStyle(
        "Body",
        parent=styles["Normal"],
        fontSize=11,
        leading=15,
        spaceAfter=6,
    )
    meta_style = ParagraphStyle(
        "Meta",
        parent=styles["Normal"],
        fontSize=9,
        textColor="#555555",
        spaceAfter=12,
    )

    story = []

    # Title
    story.append(Paragraph(_esc(lesson.title), title_style))
    story.append(
        Paragraph(
            _esc(
                f"Teacher: {persona.name or 'Teacher'}  |  "
                f"Lesson {lesson.lesson_number}  |  "
                f"{date.today().strftime('%B %d, %Y')}"
            ),
            meta_style,
        )
    )

    # Objective
    story.append(Paragraph("Objective (SWBAT)", heading_style))
    story.append(Paragraph(_esc(lesson.objective), body_style))

    # Standards
    if lesson.standards:
        story.append(Paragraph("Standards", heading_style))
        for s in lesson.standards:
            story.append(Paragraph(f"- {_esc(s)}", body_style))

    # Sections
    sections = [
        ("Do Now / Warm-Up", lesson.do_now),
        ("Direct Instruction", lesson.direct_instruction),
        ("Guided Practice", lesson.guided_practice),
        ("Independent Work", lesson.independent_work),
    ]
    for heading, content in sections:
        if content:
            story.append(Paragraph(_esc(heading), heading_style))
            story.append(Paragraph(_esc(content), body_style))

    # Exit Ticket
    if lesson.exit_ticket:
        story.append(Paragraph("Exit Ticket", heading_style))
        for i, q in enumerate(lesson.exit_ticket, 1):
            story.append(Paragraph(f"{i}. {_esc(q.question)}", body_style))

    # Homework
    if lesson.homework:
        story.append(Paragraph("Homework", heading_style))
        story.append(Paragraph(_esc(lesson.homework), body_style))

    # Footer
    story.append(Spacer(1, 24))
    footer_style = ParagraphStyle(
        "Footer", parent=styles["Normal"], fontSize=8, textColor="#888888"
    )
    story.append(Paragraph("Generated by EDUagent", footer_style))

    doc.build(story)
    return out


# ── Student handout export ─────────────────────────────────────────────


def export_student_handout(
    lesson: "DailyLesson",
    persona: "TeacherPersona",
    output_dir: Path | None = None,
) -> Path:
    """Generate a print-ready student handout (DOCX worksheet).

    Produces a clean, self-contained 1-2 page document suitable for
    printing on standard letter paper.  Includes:
    - Header with lesson title, teacher name, date
    - Bordered Do Now box with lined response area
    - Aim / Objective line
    - Core content section (direct instruction text)
    - Numbered activity questions with lined response areas
    - Exit ticket questions with response lines
    - Footer with Name / Date / Period blanks

    Returns the path to the saved .docx file.
    """
    from docx import Document
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Inches, Pt, RGBColor

    doc = Document()

    # ── Page setup ────────────────────────────────────────────────────
    for section in doc.sections:
        section.top_margin = Inches(0.6)
        section.bottom_margin = Inches(0.5)
        section.left_margin = Inches(0.75)
        section.right_margin = Inches(0.75)

    # ── Default font ──────────────────────────────────────────────────
    style = doc.styles["Normal"]
    font = style.font
    font.name = "Calibri"
    font.size = Pt(12)

    # ── Header: Title, Teacher, Date, Period ──────────────────────────
    title_para = doc.add_paragraph()
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title_para.add_run(lesson.title)
    run.bold = True
    run.font.size = Pt(14)
    run.font.name = "Calibri"

    meta_para = doc.add_paragraph()
    meta_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta_run = meta_para.add_run(
        f"{persona.name or 'Teacher'}  |  {date.today().strftime('%B %d, %Y')}"
    )
    meta_run.font.size = Pt(10)
    meta_run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    meta_run.font.name = "Calibri"

    # ── Header image (below title, max 2 images total) ───────────────
    subject = persona.subject_area or ""
    _handout_image_count = 0
    if _handout_image_count < 2:
        if _docx_add_content_image(
            doc,
            content_text=lesson.title + " " + lesson.objective,
            fallback_topic=lesson.title,
            subject=subject,
            width_inches=3.0,
        ):
            _handout_image_count += 1

    # ── Do Now box ────────────────────────────────────────────────────
    if lesson.do_now:
        _handout_section_heading(doc, "Do Now")
        do_now_table = doc.add_table(rows=2, cols=1)
        do_now_table.alignment = WD_TABLE_ALIGNMENT.CENTER
        _set_table_borders(do_now_table)

        # Prompt cell
        prompt_cell = do_now_table.rows[0].cells[0]
        prompt_cell.text = ""
        prompt_para = prompt_cell.paragraphs[0]
        prompt_run = prompt_para.add_run(lesson.do_now)
        prompt_run.font.size = Pt(11)
        prompt_run.font.name = "Calibri"

        # Response area with lines
        response_cell = do_now_table.rows[1].cells[0]
        response_cell.text = ""
        _add_lined_space(response_cell, line_count=4)

    # ── Aim / Objective ───────────────────────────────────────────────
    _handout_section_heading(doc, "Aim")
    aim_para = doc.add_paragraph(lesson.objective)
    aim_para.paragraph_format.space_after = Pt(4)

    # ── Core Content ──────────────────────────────────────────────────
    if lesson.direct_instruction:
        _handout_section_heading(doc, "Key Content")
        # Include a condensed version -- long DI sections get trimmed
        content_text = lesson.direct_instruction
        if len(content_text) > 1200:
            # Take first ~1200 chars at a sentence boundary
            cutoff = content_text[:1200].rfind(". ")
            if cutoff > 600:
                content_text = content_text[: cutoff + 1]
        content_para = doc.add_paragraph(content_text)
        content_para.paragraph_format.space_after = Pt(6)
        for run in content_para.runs:
            run.font.size = Pt(11)
            run.font.name = "Calibri"

        # Add content image next to direct instruction (max 2 total)
        if _handout_image_count < 2:
            if _docx_add_content_image(
                doc,
                content_text=lesson.direct_instruction,
                fallback_topic=lesson.title,
                subject=subject,
                width_inches=3.0,
            ):
                _handout_image_count += 1

    # ── Activity Section (Guided Practice) ────────────────────────────
    if lesson.guided_practice:
        _handout_section_heading(doc, "Activity")
        _add_numbered_content_with_lines(doc, lesson.guided_practice)

    # ── Independent Work ──────────────────────────────────────────────
    if lesson.independent_work:
        _handout_section_heading(doc, "Independent Practice")
        _add_numbered_content_with_lines(doc, lesson.independent_work)

    # ── Exit Ticket ───────────────────────────────────────────────────
    if lesson.exit_ticket:
        _handout_section_heading(doc, "Exit Ticket")
        exit_table = doc.add_table(rows=len(lesson.exit_ticket) * 2, cols=1)
        exit_table.alignment = WD_TABLE_ALIGNMENT.CENTER
        _set_table_borders(exit_table)

        for i, q in enumerate(lesson.exit_ticket):
            q_cell = exit_table.rows[i * 2].cells[0]
            q_cell.text = ""
            q_para = q_cell.paragraphs[0]
            q_run = q_para.add_run(f"{i + 1}. {q.question}")
            q_run.bold = True
            q_run.font.size = Pt(11)
            q_run.font.name = "Calibri"

            ans_cell = exit_table.rows[i * 2 + 1].cells[0]
            ans_cell.text = ""
            _add_lined_space(ans_cell, line_count=3)

    # ── Footer: Name / Date / Period ──────────────────────────────────
    doc.add_paragraph("")  # spacer
    footer_table = doc.add_table(rows=1, cols=3)
    footer_table.alignment = WD_TABLE_ALIGNMENT.CENTER
    footer_table.columns[0].width = Inches(3.0)
    footer_table.columns[1].width = Inches(2.5)
    footer_table.columns[2].width = Inches(1.5)

    labels = ["Name: _______________", "Date: _______________", "Period: ______"]
    for idx, label in enumerate(labels):
        cell = footer_table.rows[0].cells[idx]
        cell.text = ""
        para = cell.paragraphs[0]
        run = para.add_run(label)
        run.font.size = Pt(10)
        run.font.name = "Calibri"
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Remove borders from footer table
    _remove_table_borders(footer_table)

    # ── Watermark ─────────────────────────────────────────────────────
    wm_para = doc.add_paragraph()
    wm_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    wm_run = wm_para.add_run("Generated by EDUagent")
    wm_run.font.size = Pt(8)
    wm_run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    # ── Save ──────────────────────────────────────────────────────────
    out = _resolve_output(output_dir, lesson, "_handout.docx")
    doc.save(str(out))
    return out


def _handout_section_heading(doc: "Any", text: str) -> None:
    """Add a styled section heading for the student handout."""
    from docx.shared import Pt, RGBColor

    para = doc.add_paragraph()
    para.paragraph_format.space_before = Pt(8)
    para.paragraph_format.space_after = Pt(4)
    run = para.add_run(text)
    run.bold = True
    run.font.size = Pt(12)
    run.font.name = "Calibri"
    run.font.color.rgb = RGBColor(0x1A, 0x36, 0x5D)
    # Add a bottom border via XML
    from docx.oxml.ns import qn

    pPr = para._p.get_or_add_pPr()
    pBdr = pPr.makeelement(qn("w:pBdr"), {})
    bottom = pBdr.makeelement(
        qn("w:bottom"),
        {
            qn("w:val"): "single",
            qn("w:sz"): "6",
            qn("w:space"): "1",
            qn("w:color"): "1A365D",
        },
    )
    pBdr.append(bottom)
    pPr.append(pBdr)


def _add_lined_space(cell: "Any", line_count: int = 3) -> None:
    """Add blank lines with bottom borders to simulate writing lines."""
    from docx.oxml.ns import qn
    from docx.shared import Pt

    for i in range(line_count):
        para = cell.add_paragraph("")
        para.paragraph_format.space_before = Pt(0)
        para.paragraph_format.space_after = Pt(12)
        # Add bottom border to simulate a writing line
        pPr = para._p.get_or_add_pPr()
        pBdr = pPr.makeelement(qn("w:pBdr"), {})
        bottom = pBdr.makeelement(
            qn("w:bottom"),
            {
                qn("w:val"): "single",
                qn("w:sz"): "4",
                qn("w:space"): "1",
                qn("w:color"): "CCCCCC",
            },
        )
        pBdr.append(bottom)
        pPr.append(pBdr)


def _add_numbered_content_with_lines(doc: "Any", text: str) -> None:
    """Parse text into numbered items and add lined response areas.

    If the text contains numbered items (1. / 2. / etc.), each gets its own
    response area.  Otherwise the full text is shown with a single response area.
    """
    import re

    from docx.shared import Pt

    # Try to split on numbered items
    items = re.split(r"(?:^|\n)\s*(\d+)[.)]\s*", text)

    # items[0] is preamble, then alternating [number, content, number, content...]
    preamble = items[0].strip() if items else ""
    numbered: list[str] = []
    if len(items) > 2:
        for i in range(1, len(items) - 1, 2):
            numbered.append(items[i + 1].strip() if i + 1 < len(items) else "")

    if numbered:
        if preamble:
            p = doc.add_paragraph(preamble)
            p.paragraph_format.space_after = Pt(4)
            for run in p.runs:
                run.font.size = Pt(11)
                run.font.name = "Calibri"
        for idx, item in enumerate(numbered, 1):
            p = doc.add_paragraph()
            run = p.add_run(f"{idx}. {item}")
            run.font.size = Pt(11)
            run.font.name = "Calibri"
            p.paragraph_format.space_after = Pt(2)
            # Add a small lined area after each item
            line_para = doc.add_paragraph("")
            line_para.paragraph_format.space_after = Pt(10)
            from docx.oxml.ns import qn

            pPr = line_para._p.get_or_add_pPr()
            pBdr = pPr.makeelement(qn("w:pBdr"), {})
            bottom = pBdr.makeelement(
                qn("w:bottom"),
                {
                    qn("w:val"): "single",
                    qn("w:sz"): "4",
                    qn("w:space"): "1",
                    qn("w:color"): "CCCCCC",
                },
            )
            pBdr.append(bottom)
            pPr.append(pBdr)
    else:
        # No numbered items — show the text with a response area
        p = doc.add_paragraph(text)
        p.paragraph_format.space_after = Pt(4)
        for run in p.runs:
            run.font.size = Pt(11)
            run.font.name = "Calibri"
        # Add 3 lined spaces
        for _ in range(3):
            line_para = doc.add_paragraph("")
            line_para.paragraph_format.space_after = Pt(10)
            from docx.oxml.ns import qn

            pPr = line_para._p.get_or_add_pPr()
            pBdr = pPr.makeelement(qn("w:pBdr"), {})
            bottom = pBdr.makeelement(
                qn("w:bottom"),
                {
                    qn("w:val"): "single",
                    qn("w:sz"): "4",
                    qn("w:space"): "1",
                    qn("w:color"): "CCCCCC",
                },
            )
            pBdr.append(bottom)
            pPr.append(pBdr)


def _set_table_borders(table: "Any") -> None:
    """Set solid borders on all sides of a table."""
    from docx.oxml.ns import qn

    tbl = table._tbl
    tblPr = tbl.tblPr if tbl.tblPr is not None else tbl.makeelement(qn("w:tblPr"), {})
    borders = tblPr.makeelement(qn("w:tblBorders"), {})
    for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
        element = borders.makeelement(
            qn(f"w:{edge}"),
            {
                qn("w:val"): "single",
                qn("w:sz"): "6",
                qn("w:space"): "0",
                qn("w:color"): "444444",
            },
        )
        borders.append(element)
    tblPr.append(borders)


def _remove_table_borders(table: "Any") -> None:
    """Remove all borders from a table."""
    from docx.oxml.ns import qn

    tbl = table._tbl
    tblPr = tbl.tblPr if tbl.tblPr is not None else tbl.makeelement(qn("w:tblPr"), {})
    borders = tblPr.makeelement(qn("w:tblBorders"), {})
    for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
        element = borders.makeelement(
            qn(f"w:{edge}"),
            {qn("w:val"): "none", qn("w:sz"): "0", qn("w:space"): "0"},
        )
        borders.append(element)
    tblPr.append(borders)


# ── Helpers ────────────────────────────────────────────────────────────


def _resolve_output(output_dir: Path | None, lesson: "DailyLesson", ext: str) -> Path:
    """Build the output file path."""
    if output_dir is None:
        output_dir = Path("eduagent_output").resolve()
    output_dir.mkdir(parents=True, exist_ok=True)
    safe = f"lesson_{lesson.lesson_number:02d}"
    return output_dir / f"{safe}{ext}"


def _split_text(text: str, max_len: int = 550) -> list[str]:
    """Split long text into chunks at sentence boundaries."""
    sentences = text.replace("\n", " ").split(". ")
    chunks: list[str] = []
    current = ""
    for s in sentences:
        candidate = f"{current}. {s}" if current else s
        if len(candidate) > max_len and current:
            chunks.append(current.strip())
            current = s
        else:
            current = candidate
    if current.strip():
        chunks.append(current.strip())
    return chunks or [text]


def _esc(text: str) -> str:
    """Escape text for reportlab XML-based Paragraphs."""
    return (
        text.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
    )
