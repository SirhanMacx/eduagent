"""File ingestion pipeline — extracts text from teacher curriculum files.

Supported formats: PDF, DOCX, PPTX, TXT, MD, SMART Notebook (.notebook),
SMART Board (.xbk), ActivInspire (.flipchart), and ZIP archives.
"""

from __future__ import annotations

import logging
import tempfile
import zipfile
from collections import Counter
from pathlib import Path
from typing import Optional

from eduagent.models import DocType, Document

logger = logging.getLogger(__name__)

# ── File-type detection ──────────────────────────────────────────────────

EXTENSION_MAP: dict[str, DocType] = {
    ".pdf": DocType.PDF,
    ".docx": DocType.DOCX,
    ".pptx": DocType.PPTX,
    ".txt": DocType.TXT,
    ".md": DocType.MD,
    ".notebook": DocType.NOTEBOOK,
    ".xbk": DocType.XBK,
    ".flipchart": DocType.FLIPCHART,
}

SUPPORTED_EXTENSIONS = set(EXTENSION_MAP.keys())


def _detect_type(path: Path) -> DocType:
    return EXTENSION_MAP.get(path.suffix.lower(), DocType.UNKNOWN)


# ── Individual extractors ────────────────────────────────────────────────


def _extract_pdf(path: Path) -> tuple[str, Optional[int]]:
    """Extract text from a PDF using PyMuPDF."""
    import fitz  # PyMuPDF

    doc = fitz.open(str(path))
    pages: list[str] = []
    for page in doc:
        text = page.get_text()
        if text.strip():
            pages.append(text)
    page_count = len(doc)
    doc.close()
    return "\n\n".join(pages), page_count


def _extract_docx(path: Path) -> str:
    """Extract text from a DOCX file."""
    from docx import Document as DocxDocument

    doc = DocxDocument(str(path))
    paragraphs: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            paragraphs.append(para.text)
    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))
    return "\n\n".join(paragraphs)


def _extract_pptx(path: Path) -> str:
    """Extract text from a PPTX file."""
    from pptx import Presentation

    prs = Presentation(str(path))
    slides_text: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts: list[str] = [f"[Slide {i}]"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        parts.append(text)
        if len(parts) > 1:  # More than just the slide marker
            slides_text.append("\n".join(parts))
    return "\n\n".join(slides_text)


def _extract_text(path: Path) -> str:
    """Extract text from a plain text or Markdown file."""
    return path.read_text(encoding="utf-8", errors="replace")


def _extract_notebook(path: Path) -> tuple[str, Optional[int]]:
    """Extract text from a SMART Notebook file via formats plugin."""
    from eduagent.formats.notebook import extract_notebook

    return extract_notebook(path)


def _extract_xbk(path: Path) -> tuple[str, Optional[int]]:
    """Extract text from a SMART Board file via formats plugin."""
    from eduagent.formats.xbk import extract_xbk

    return extract_xbk(path)


def _extract_flipchart(path: Path) -> tuple[str, Optional[int]]:
    """Extract text from an ActivInspire file via formats plugin."""
    from eduagent.formats.flipchart import extract_flipchart

    return extract_flipchart(path)


# ── Dispatch ─────────────────────────────────────────────────────────────

EXTRACTORS = {
    DocType.PDF: lambda p: _extract_pdf(p),
    DocType.DOCX: lambda p: (_extract_docx(p), None),
    DocType.PPTX: lambda p: (_extract_pptx(p), None),
    DocType.TXT: lambda p: (_extract_text(p), None),
    DocType.MD: lambda p: (_extract_text(p), None),
    DocType.NOTEBOOK: lambda p: _extract_notebook(p),
    DocType.XBK: lambda p: _extract_xbk(p),
    DocType.FLIPCHART: lambda p: _extract_flipchart(p),
}


def _extract_single(path: Path) -> Optional[Document]:
    """Extract a single file into a Document, or None if unsupported/empty."""
    doc_type = _detect_type(path)
    if doc_type == DocType.UNKNOWN:
        logger.warning("Skipping unsupported format: %s", path.name)
        return None

    extractor = EXTRACTORS.get(doc_type)
    if not extractor:
        logger.warning("No extractor for %s format: %s", doc_type.value, path.name)
        return None

    try:
        result = extractor(path)
    except Exception as exc:
        logger.warning("Failed to parse %s: %s", path.name, exc)
        return None
    if isinstance(result, tuple):
        content, page_count = result
    else:
        content, page_count = result, None

    if not content or not content.strip():
        return None

    return Document(
        title=path.stem.replace("_", " ").replace("-", " ").title(),
        content=content.strip(),
        doc_type=doc_type,
        source_path=str(path),
        page_count=page_count,
    )


# ── Format summary ──────────────────────────────────────────────────────


def _collect_files(path: Path) -> list[Path]:
    """Collect all supported files from a directory, sorted.

    Skips Office lock files (~$*) and macOS resource forks (._*).
    """
    return sorted(
        f for f in path.rglob("*")
        if f.is_file()
        and f.suffix.lower() in SUPPORTED_EXTENSIONS
        and not f.name.startswith(("~$", "._"))
    )


def _format_summary(files: list[Path]) -> str:
    """Build a human-readable summary of found file types.

    Example: "Found 47 PDFs, 23 DOCX, 12 PPTX, 3 TXT — analyzing..."
    """
    counts: Counter[str] = Counter()
    for f in files:
        ext = f.suffix.lower().lstrip(".")
        counts[ext.upper()] += 1

    parts = [f"{count} {fmt}" for fmt, count in sorted(counts.items(), key=lambda x: -x[1])]
    return f"Found {', '.join(parts)} — analyzing..."


# ── Public API ───────────────────────────────────────────────────────────


def scan_directory(path: Path) -> tuple[list[Path], str]:
    """Scan a directory and return supported files with a format summary.

    Returns (file_list, summary_string).
    """
    if not path.is_dir():
        raise FileNotFoundError(f"Directory not found: {path}")

    files = _collect_files(path)
    summary = _format_summary(files) if files else "No supported files found."
    return files, summary


def ingest_directory(path: Path, *, progress_callback=None) -> list[Document]:
    """Recursively scan a directory and extract all supported documents.

    Args:
        path: Directory to scan.
        progress_callback: Optional callable(current, total) for progress updates.
    """
    documents: list[Document] = []
    if not path.is_dir():
        raise FileNotFoundError(f"Directory not found: {path}")

    files = _collect_files(path)
    total = len(files)

    for i, file_path in enumerate(files):
        doc = _extract_single(file_path)
        if doc:
            documents.append(doc)
        if progress_callback:
            progress_callback(i + 1, total)

    return documents


def ingest_zip(path: Path, *, progress_callback=None) -> list[Document]:
    """Extract a ZIP file to a temp directory, then ingest its contents."""
    if not path.is_file() or path.suffix.lower() != ".zip":
        raise ValueError(f"Not a ZIP file: {path}")

    with tempfile.TemporaryDirectory() as tmp_dir:
        with zipfile.ZipFile(str(path), "r") as zf:
            zf.extractall(tmp_dir)
        return ingest_directory(Path(tmp_dir), progress_callback=progress_callback)


def ingest_path(path: Path, *, dry_run: bool = False, progress_callback=None) -> list[Document]:
    """Smart ingestion: accepts a directory, ZIP file, or single file.

    Args:
        path: File or directory to ingest.
        dry_run: If True, scan and log what would be processed but don't extract.
        progress_callback: Optional callable(current, total) for progress updates.
    """
    path = Path(path).expanduser().resolve()

    if path.is_dir():
        if dry_run:
            files, summary = scan_directory(path)
            logger.info(summary)
            return _dry_run_results(files)
        return ingest_directory(path, progress_callback=progress_callback)
    elif path.is_file() and path.suffix.lower() == ".zip":
        if dry_run:
            # For ZIP dry-run, list contents without extracting
            with tempfile.TemporaryDirectory() as tmp_dir:
                with zipfile.ZipFile(str(path), "r") as zf:
                    zf.extractall(tmp_dir)
                files, summary = scan_directory(Path(tmp_dir))
                logger.info(summary)
                return _dry_run_results(files)
        return ingest_zip(path, progress_callback=progress_callback)
    elif path.is_file():
        if dry_run:
            return _dry_run_results([path])
        doc = _extract_single(path)
        return [doc] if doc else []
    else:
        raise FileNotFoundError(f"Path not found: {path}")


def _dry_run_results(files: list[Path]) -> list[Document]:
    """Create placeholder Documents for dry-run mode (no content extraction)."""
    results: list[Document] = []
    for f in files:
        doc_type = _detect_type(f)
        if doc_type != DocType.UNKNOWN:
            results.append(Document(
                title=f.stem.replace("_", " ").replace("-", " ").title(),
                content="[dry-run: content not extracted]",
                doc_type=doc_type,
                source_path=str(f),
            ))
    return results
